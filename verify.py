"""
Quality gates. A deck that fails any check does not ship.

    python verify.py                 # all decks in output/
    python verify.py --pp            # also open every deck in PowerPoint
    python verify.py --urls          # also check every external URL

What each check measures is stated in the report, so nothing is hidden behind
a pass. Text overflow is measured, not eyeballed: every text frame in this
library has zero internal margin and explicit line spacing, so predicted
height is exact.
"""

import glob
import json
import os
import re
import sys
import zipfile

from lxml import etree
from pptx import Presentation
from pptx.util import Pt

import theme as T
import textfit as TF

ROOT = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(ROOT, "output")
EMU = 914400.0

BANNED = ["lorem", "todo", "placeholder", "tbd", "insert here", "xxx",
          "leverage", "synergy", "in today's fast-paced world",
          "unlock the power of", "coming soon"]

# Prose blocks — the "no bullet wall" rule applies to these.
PROSE = "body:"
# Structured component cells — capped individually instead.
CELL = "cell:"

MAX_PROSE_LINES = 5
MAX_PROSE_WORDS = 60
MAX_CELL_WORDS = 34
MAX_SLIDE_WORDS = 175
MAX_INK_COVERAGE = 0.70          # i.e. at least 30% of the slide stays empty
OVERFLOW_TOL = 0.03              # inches


class Report:
    def __init__(self, name):
        self.name = name
        self.fails = []
        self.warns = []
        self.stats = {}

    def fail(self, check, where, msg):
        self.fails.append((check, where, msg))

    def warn(self, check, where, msg):
        self.warns.append((check, where, msg))

    @property
    def ok(self):
        return not self.fails


# --------------------------------------------------------------------------
# shape helpers
# --------------------------------------------------------------------------
def box(sh):
    return (sh.left / EMU, sh.top / EMU, sh.width / EMU, sh.height / EMU)


def solid_fill_rgb(sh):
    try:
        if sh.fill.type is not None and sh.fill.type == 1:      # MSO_FILL.SOLID
            return sh.fill.fore_color.rgb
    except Exception:
        pass
    return None


def runs_of(sh):
    out = []
    if not sh.has_text_frame:
        return out
    for para in sh.text_frame.paragraphs:
        ls = para.line_spacing
        ls_pt = ls.pt if hasattr(ls, "pt") else None
        sa = para.space_after
        sa_pt = sa.pt if hasattr(sa, "pt") else 0.0
        for run in para.runs:
            if not run.text.strip():
                continue
            size = run.font.size.pt if run.font.size else None
            name = run.font.name or T.F_BODY
            col = None
            try:
                if run.font.color and run.font.color.type is not None:
                    col = run.font.color.rgb
            except Exception:
                col = None
            out.append({"text": run.text, "pt": size, "font": name,
                        "color": col, "line_pt": ls_pt, "sa_pt": sa_pt,
                        "para": para})
    return out


def slide_bg(slide):
    for sh in slide.shapes:
        if sh.name == "vis:bg":
            c = solid_fill_rgb(sh)
            if c is not None:
                return c
    return T.SURFACE


def bg_under(slide, target, upto_index):
    """
    Topmost solid-filled shape that fully contains `target` and sits below it
    in z-order. Falls back to the slide background.
    """
    x, y, w, h = box(target)
    found = slide_bg(slide)
    for i, sh in enumerate(slide.shapes):
        if i >= upto_index:
            break
        c = solid_fill_rgb(sh)
        if c is None:
            continue
        sx, sy, sw, sh_ = box(sh)
        if sx - 0.02 <= x and sy - 0.02 <= y and \
           sx + sw + 0.02 >= x + w and sy + sh_ + 0.02 >= y + h:
            found = c
    return found


# --------------------------------------------------------------------------
# individual checks
# --------------------------------------------------------------------------
def check_opens(path, rep):
    """Gate 1 — the package is well-formed and every slide part parses."""
    try:
        prs = Presentation(path)
    except Exception as exc:
        rep.fail("1 opens", path, "python-pptx could not open: %s" % exc)
        return None
    try:
        with zipfile.ZipFile(path) as z:
            for n in z.namelist():
                if n.endswith(".xml") or n.endswith(".rels"):
                    etree.fromstring(z.read(n))
    except Exception as exc:
        rep.fail("1 opens", path, "malformed XML part: %s" % exc)
    return prs


def check_slide(slide, idx, rep):
    shapes = list(slide.shapes)
    visual = 0
    slide_words = 0

    for i, sh in enumerate(shapes):
        role = sh.name.split(":")[0] if ":" in sh.name else sh.name
        is_txtbox = sh.shape_type is not None and str(sh.shape_type).startswith("TEXT_BOX")

        if sh.name == "vis:bg":
            pass
        elif solid_fill_rgb(sh) is not None or not sh.has_text_frame:
            visual += 1
        elif sh.shape_type is not None and "PICTURE" in str(sh.shape_type):
            visual += 1

        rs = runs_of(sh)
        if not rs:
            continue

        x, y, w, h = box(sh)
        tf = sh.text_frame
        ml = (tf.margin_left or 0) / EMU
        mr = (tf.margin_right or 0) / EMU
        mt = (tf.margin_top or 0) / EMU
        mb = (tf.margin_bottom or 0) / EMU
        avail_w = w - ml - mr
        avail_h = h - mt - mb

        # --- gate 6: minimum font size
        for r in rs:
            if r["pt"] is None:
                rep.fail("6 font size", "slide %d" % idx,
                         "%s: run has no explicit size" % sh.name)
            elif r["pt"] < T.MIN_PT - 0.01:
                rep.fail("6 font size", "slide %d" % idx,
                         "%s: %.1fpt < %dpt  (%r)" %
                         (sh.name, r["pt"], T.MIN_PT, r["text"][:40]))

        # --- gate 2: text must fit its shape
        total_h = 0.0
        widest = 0.0
        lines_total = 0
        for r in rs:
            pt = r["pt"] or T.SZ_BODY
            lp = r["line_pt"] or pt * TF.LINE_FACTOR
            lines = TF.wrap(r["text"], r["font"], pt, avail_w)
            lines_total += len(lines)
            total_h += len(lines) * lp / 72.0 + (r["sa_pt"] or 0.0) / 72.0
            for ln in lines:
                widest = max(widest, TF.text_w_in(ln, r["font"], pt))
        if total_h > avail_h + OVERFLOW_TOL:
            rep.fail("2 overflow", "slide %d" % idx,
                     "%s: text needs %.2fin, shape gives %.2fin  (%r)" %
                     (sh.name, total_h, avail_h, rs[0]["text"][:48]))
        if widest > avail_w + OVERFLOW_TOL:
            rep.fail("2 overflow", "slide %d" % idx,
                     "%s: longest line %.2fin > width %.2fin" %
                     (sh.name, widest, avail_w))

        # --- design rule: titles stay on one line at the 34pt scale
        if sh.name == "title" and lines_total > 1 and (rs[0]["pt"] or 0) >= 30:
            rep.warn("title", "slide %d" % idx,
                     "title wraps to %d lines" % lines_total)
        if sh.name == "title" and (rs[0]["pt"] or 0) < T.SZ_TITLE - 0.01 \
                and (rs[0]["pt"] or 0) >= 20:
            rep.warn("title", "slide %d" % idx,
                     "title shrank to %.1fpt (scale is %dpt) — shorten the text"
                     % (rs[0]["pt"], T.SZ_TITLE))

        # --- gate 7: contrast
        bg = bg_under(slide, sh, i)
        for r in rs:
            if r["color"] is None:
                continue
            ratio = T.contrast(r["color"], bg)
            if ratio < T.AA:
                rep.fail("7 contrast", "slide %d" % idx,
                         "%s: %.2f:1  #%s on #%s  (%r)" %
                         (sh.name, ratio, r["color"], bg, r["text"][:36]))

        words = sum(len(r["text"].split()) for r in rs)
        # A copy-paste prompt is a template to copy, not text to read, so it
        # does not count toward the slide's reading load.
        if sh.name != "card:prompt":
            slide_words += words

        # --- gate 5: no bullet walls in prose blocks
        if sh.name.startswith(PROSE):
            if lines_total > MAX_PROSE_LINES:
                rep.fail("5 density", "slide %d" % idx,
                         "%s: %d body lines > %d" %
                         (sh.name, lines_total, MAX_PROSE_LINES))
            if words > MAX_PROSE_WORDS:
                rep.fail("5 density", "slide %d" % idx,
                         "%s: %d body words > %d" %
                         (sh.name, words, MAX_PROSE_WORDS))
        elif sh.name.startswith(CELL):
            if words > MAX_CELL_WORDS:
                rep.fail("5 density", "slide %d" % idx,
                         "%s: component cell has %d words > %d" %
                         (sh.name, words, MAX_CELL_WORDS))

        # --- gate 9: banned strings
        low = " ".join(r["text"] for r in rs).lower()
        for bad in BANNED:
            if bad in low:
                rep.fail("9 banned", "slide %d" % idx,
                         "%s: contains %r" % (sh.name, bad))

    # --- gate 8: every slide carries a visual
    if visual == 0:
        rep.fail("8 visual", "slide %d" % idx, "slide is text-only")

    if slide_words > MAX_SLIDE_WORDS:
        rep.fail("5 density", "slide %d" % idx,
                 "%d words on one slide > %d" % (slide_words, MAX_SLIDE_WORDS))

    # --- gate 11: whitespace
    cov = ink_coverage(slide)
    if cov > MAX_INK_COVERAGE:
        rep.fail("11 whitespace", "slide %d" % idx,
                 "%.0f%% of the slide is covered (max %.0f%%)" %
                 (cov * 100, MAX_INK_COVERAGE * 100))
    return slide_words


def text_rect(sh):
    """Rendered extent of a text shape: where the glyphs actually land."""
    rs = runs_of(sh)
    if not rs:
        return None
    x, y, w, h = box(sh)
    tf = sh.text_frame
    ml = (tf.margin_left or 0) / EMU
    mt = (tf.margin_top or 0) / EMU
    aw = w - ml - ((tf.margin_right or 0) / EMU)
    ah = h - mt - ((tf.margin_bottom or 0) / EMU)
    th = 0.0
    widest = 0.0
    for r in rs:
        pt = r["pt"] or T.SZ_BODY
        lp = r["line_pt"] or pt * TF.LINE_FACTOR
        lines = TF.wrap(r["text"], r["font"], pt, aw)
        th += len(lines) * lp / 72.0 + (r["sa_pt"] or 0.0) / 72.0
        for ln in lines:
            widest = max(widest, TF.text_w_in(ln, r["font"], pt))
    th = min(th, ah)
    widest = min(widest, aw)
    anchor = str(tf.vertical_anchor)
    ty = y + mt
    if anchor.startswith("MIDDLE"):
        ty = y + mt + (ah - th) / 2
    elif anchor.startswith("BOTTOM"):
        ty = y + mt + (ah - th)
    return (x + ml, ty, widest, th)


def check_overlap(slide, idx, rep):
    """
    Gate 2b - no two blocks of text may sit on top of each other.
    Text over a filled card is normal; text over text never is.
    """
    rects = []
    for sh in slide.shapes:
        r = text_rect(sh)
        if r and r[2] > 0.02 and r[3] > 0.02:
            rects.append((sh.name, r))
    for i in range(len(rects)):
        for j in range(i + 1, len(rects)):
            n1, (x1, y1, w1, h1) = rects[i]
            n2, (x2, y2, w2, h2) = rects[j]
            ox = min(x1 + w1, x2 + w2) - max(x1, x2)
            oy = min(y1 + h1, y2 + h2) - max(y1, y2)
            if ox > 0.04 and oy > 0.04 and ox * oy > 0.02:
                rep.fail("2b overlap", "slide %d" % idx,
                         "%s overlaps %s by %.2fin x %.2fin" %
                         (n1, n2, ox, oy))


def ink_coverage(slide, nx=140, ny=79):
    """
    Visual density: the fraction of the slide occupied by TEXT (measured at its
    rendered extent, not its box) and by STRONG blocks — fills that contrast
    2:1 or more against the slide ground.

    A pale panel (SURFACE_ALT on white, 1.08:1) reads as background, not as
    ink, so it is not counted. The full-bleed background is never counted.
    This measures crowding, which is what the whitespace rule is about.
    """
    grid = bytearray(nx * ny)
    cw = T.SLIDE_W / nx
    ch = T.SLIDE_H / ny
    ground = slide_bg(slide)

    def mark(x, y, w, h):
        i0 = max(0, int(x / cw)); i1 = min(nx, int((x + w) / cw) + 1)
        j0 = max(0, int(y / ch)); j1 = min(ny, int((y + h) / ch) + 1)
        for j in range(j0, j1):
            row = j * nx
            for i in range(i0, i1):
                grid[row + i] = 1

    for sh in slide.shapes:
        if sh.name == "vis:bg":
            continue
        x, y, w, h = box(sh)
        fill = solid_fill_rgb(sh)
        if fill is not None and T.contrast(fill, ground) >= 2.0:
            mark(x, y, w, h)
            continue
        rs = runs_of(sh)
        if not rs:
            continue
        tf = sh.text_frame
        ml = (tf.margin_left or 0) / EMU
        mt = (tf.margin_top or 0) / EMU
        aw = w - ml - ((tf.margin_right or 0) / EMU)
        ah = h - mt - ((tf.margin_bottom or 0) / EMU)
        th = 0.0
        widest = 0.0
        for r in rs:
            pt = r["pt"] or T.SZ_BODY
            lp = r["line_pt"] or pt * TF.LINE_FACTOR
            lines = TF.wrap(r["text"], r["font"], pt, aw)
            th += len(lines) * lp / 72.0 + (r["sa_pt"] or 0.0) / 72.0
            for ln in lines:
                widest = max(widest, TF.text_w_in(ln, r["font"], pt))
        th = min(th, ah)
        widest = min(widest, aw)
        anchor = tf.vertical_anchor
        ty = y + mt
        if str(anchor).startswith("MIDDLE"):
            ty = y + mt + (ah - th) / 2
        elif str(anchor).startswith("BOTTOM"):
            ty = y + mt + (ah - th)
        mark(x + ml, ty, widest, th)
    return sum(grid) / float(nx * ny)


def check_links(path, prs, rep):
    """Gate 3 — every internal jump lands on a slide that exists."""
    with zipfile.ZipFile(path) as z:
        names = set(z.namelist())
        slide_parts = {n for n in names if
                       re.match(r"ppt/slides/slide\d+\.xml$", n)}
        total = 0
        for n in sorted(slide_parts):
            rels = "ppt/slides/_rels/%s.rels" % os.path.basename(n)
            if rels not in names:
                continue
            root = etree.fromstring(z.read(rels))
            targets = {}
            for rel in root:
                targets[rel.get("Id")] = (rel.get("Type"), rel.get("Target"))
            sx = etree.fromstring(z.read(n))
            ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main",
                  "r": "http://schemas.openxmlformats.org/officeDocument/"
                       "2006/relationships"}
            for hl in sx.iter("{%s}hlinkClick" % ns["a"]):
                rid = hl.get("{%s}id" % ns["r"])
                action = hl.get("action") or ""
                if "jump" in action or "slide" in action:
                    total += 1
                    if not rid:
                        rep.fail("3 links", os.path.basename(n),
                                 "slide jump with no relationship id")
                        continue
                    typ, tgt = targets.get(rid, (None, None))
                    if tgt is None:
                        rep.fail("3 links", os.path.basename(n),
                                 "hyperlink rId %s has no relationship" % rid)
                        continue
                    resolved = os.path.normpath(
                        os.path.join("ppt/slides", tgt)).replace("\\", "/")
                    if resolved not in slide_parts:
                        rep.fail("3 links", os.path.basename(n),
                                 "jump target %s does not exist" % tgt)
        rep.stats["internal_jumps"] = total


def check_urls(urls, rep):
    """Gate 4 — every external URL resolves."""
    import requests
    hdr = {"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64)"}
    for u in urls:
        code, err = None, None
        # Two attempts: a slow government site timing out once is a network
        # fact, not a dead link.
        for attempt in range(2):
            try:
                r = requests.get(u, headers=hdr, timeout=40,
                                 allow_redirects=True)
                code, err = r.status_code, None
                if code == 200:
                    break
            except Exception as exc:
                err = exc
        if err is not None:
            rep.fail("4 urls", u, "request failed twice: %s" % err)
        elif code != 200:
            rep.fail("4 urls", u, "HTTP %s" % code)


def check_powerpoint(path, rep):
    """Gate 1b — PowerPoint itself opens the file with no repair prompt."""
    try:
        import win32com.client
        app = win32com.client.Dispatch("PowerPoint.Application")
        app.DisplayAlerts = 1
        pres = app.Presentations.Open(os.path.abspath(path), ReadOnly=True,
                                      WithWindow=False)
        n = pres.Slides.Count
        pres.Close()
        app.Quit()
        rep.stats["powerpoint_slides"] = n
    except Exception as exc:
        rep.fail("1b powerpoint", path, "PowerPoint could not open: %s" % exc)


# --------------------------------------------------------------------------
def verify_deck(path, manifest_entry=None, do_urls=False, do_pp=False):
    rep = Report(os.path.basename(path))
    prs = check_opens(path, rep)
    if prs is None:
        return rep

    total_words = 0
    for idx, slide in enumerate(prs.slides, start=1):
        total_words += check_slide(slide, idx, rep) or 0
        check_overlap(slide, idx, rep)
    rep.stats["slides"] = len(prs.slides._sldIdLst)
    rep.stats["words"] = total_words

    check_links(path, prs, rep)

    if manifest_entry:
        cp = (manifest_entry.get("content_points") or {}).get("total", 0)
        rep.stats["content_points"] = cp
        if not manifest_entry.get("is_index") and not (25 <= cp <= 40):
            rep.fail("10 content points", rep.name,
                     "%d points, must be 25-40" % cp)
        if do_urls:
            check_urls(manifest_entry.get("external_links", []), rep)
    if do_pp:
        check_powerpoint(path, rep)
    return rep


def main(argv):
    try:
        sys.stdout.reconfigure(encoding="utf-8")
    except Exception:
        pass
    do_urls = "--urls" in argv
    do_pp = "--pp" in argv
    only = [a for a in argv[1:] if not a.startswith("--")]

    mpath = os.path.join(OUT, "manifest.json")
    manifest = {}
    if os.path.exists(mpath):
        with open(mpath, encoding="utf-8") as fh:
            manifest = {d["file"]: d for d in json.load(fh).get("decks", [])}

    files = sorted(glob.glob(os.path.join(OUT, "*", "*.pptx"))) + \
        sorted(glob.glob(os.path.join(OUT, "*.pptx")))
    if only:
        files = [f for f in files if any(o.lower() in f.lower() for o in only)]

    reports = []
    for f in files:
        rel = os.path.relpath(f, ROOT).replace("\\", "/")
        rep = verify_deck(f, manifest.get(rel), do_urls, do_pp)
        reports.append(rep)

    print("=" * 78)
    print("INDUCTO QUALITY GATES")
    print("=" * 78)
    for rep in reports:
        st = rep.stats
        head = "%-40s %s" % (rep.name, "PASS" if rep.ok else "FAIL")
        print("\n" + head)
        print("  slides %s | content points %s | internal jumps %s | words %s" %
              (st.get("slides", "?"), st.get("content_points", "-"),
               st.get("internal_jumps", "?"), st.get("words", "?")))
        for c, w, m in rep.fails:
            print("  FAIL  [%s] %s: %s" % (c, w, m))
        for c, w, m in rep.warns:
            print("  warn  [%s] %s: %s" % (c, w, m))

    n_decks = len(files)
    print("\n" + "-" * 78)
    print("12 file count: %d topic decks + %d master index found "
          "(target 39 + 1)" %
          (len([f for f in files if "master-index" not in f]),
           len([f for f in files if "master-index" in f])))
    failed = [r.name for r in reports if not r.ok]
    print("RESULT: %d/%d decks PASS" % (n_decks - len(failed), n_decks))
    if failed:
        print("FAILED: %s" % ", ".join(failed))
    return 1 if failed else 0


if __name__ == "__main__":
    sys.exit(main(sys.argv))
