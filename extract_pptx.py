"""
Phase 1 — inspect every PPT/PPTX in output/ and record what is actually
inside it. This reads the shipped files, not the content modules, so the
inventory is evidence rather than assertion.
"""

import glob
import json
import os
import sys
from collections import Counter

from pptx import Presentation

ROOT = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(ROOT, "output")


def shape_texts(shapes):
    for sh in shapes:
        if sh.has_text_frame:
            t = sh.text_frame.text.strip()
            if t:
                yield sh.name, t


def inspect(path):
    prs = Presentation(path)
    slides = []
    roles = Counter()
    pics = 0
    for idx, sl in enumerate(prs.slides, start=1):
        hidden = sl._element.get("show") == "0"
        texts = []
        for name, t in shape_texts(sl.shapes):
            roles[name.split(":")[0] if ":" in name else name] += 1
            texts.append({"role": name, "text": t})
        for sh in sl.shapes:
            if sh.shape_type is not None and "PICTURE" in str(sh.shape_type):
                pics += 1
        title = next((t["text"] for t in texts
                      if t["role"] in ("title", "title:banner", "title:cover")), "")
        slides.append({"n": idx, "hidden": hidden, "title": title,
                       "shapes": len(sl.shapes), "texts": texts})
    words = sum(len(t["text"].split()) for s in slides for t in s["texts"])
    return {
        "file": os.path.relpath(path, ROOT).replace("\\", "/"),
        "name": os.path.basename(path),
        "slides": len(slides),
        "hidden": sum(1 for s in slides if s["hidden"]),
        "pictures": pics,
        "words": words,
        "roles": dict(roles),
        "slide_detail": slides,
    }


def main():
    sys.stdout.reconfigure(encoding="utf-8")
    files = sorted(glob.glob(os.path.join(OUT, "**", "*.ppt*"), recursive=True))
    inv = [inspect(f) for f in files]
    with open(os.path.join(ROOT, "pptx_inventory.json"), "w",
              encoding="utf-8") as fh:
        json.dump(inv, fh, indent=1, ensure_ascii=False)
    print("inspected %d files" % len(inv))
    print("slides    %d" % sum(i["slides"] for i in inv))
    print("hidden    %d" % sum(i["hidden"] for i in inv))
    print("pictures  %d" % sum(i["pictures"] for i in inv))
    print("words     %d" % sum(i["words"] for i in inv))
    return inv


if __name__ == "__main__":
    main()
