"""
Inducto visual library.

Everything here is a NATIVE PowerPoint shape so the client can edit it later.
Charts are the only thing allowed to fall back to a rendered PNG.

One line weight, one corner radius, one icon grammar — see theme.py.
"""

import copy
import hashlib
import math
import os

from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

import theme as T
import textfit as TF


# ==========================================================================
# Paragraph record
# ==========================================================================
def P(text, font=None, pt=None, color=None, sa=0, align="l", spc=0.0,
      indent=0.0, bold=False):
    return {
        "text": text,
        "font": font or T.F_BODY,
        "pt": pt or T.SZ_BODY,
        "color": color if color is not None else T.INK,
        "space_after_pt": sa,
        "align": align,
        "spc_pt": spc,
        "bullet_indent_in": indent,
        "bold": bold,
    }


_ALIGN = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}
_ANCHOR = {"t": MSO_ANCHOR.TOP, "m": MSO_ANCHOR.MIDDLE, "b": MSO_ANCHOR.BOTTOM}


# ==========================================================================
# Low-level primitives
# ==========================================================================
def _clean(sh):
    """Kill inherited theme shadow so shapes render flat and identical."""
    try:
        sh.shadow.inherit = False
    except Exception:
        pass
    return sh


def _no_fill(sh):
    sh.fill.background()
    return sh


def _no_line(sh):
    sh.line.fill.background()
    return sh


def _radius(sh, w_in, h_in, r_in=None):
    r_in = T.RADIUS if r_in is None else r_in
    small = max(0.01, min(w_in, h_in))
    try:
        sh.adjustments[0] = min(0.5, r_in / small)
    except Exception:
        pass
    return sh


def shape(sl, kind, x, y, w, h, fill=None, line=None, lw=None,
          name="vis:shape", radius=None):
    sh = sl.shapes.add_shape(kind, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.name = name
    _clean(sh)
    if kind == MSO_SHAPE.ROUNDED_RECTANGLE:
        _radius(sh, w, h, radius)
    if fill is None:
        _no_fill(sh)
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        _no_line(sh)
    else:
        sh.line.color.rgb = line
        sh.line.width = lw or T.LINE_W
    sh.text_frame.text = ""
    sh.text_frame.word_wrap = True
    return sh


def rect(sl, x, y, w, h, **kw):
    kw.setdefault("name", "vis:rect")
    return shape(sl, MSO_SHAPE.RECTANGLE, x, y, w, h, **kw)


def rrect(sl, x, y, w, h, **kw):
    kw.setdefault("name", "vis:card")
    return shape(sl, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, **kw)


def oval(sl, x, y, w, h, **kw):
    kw.setdefault("name", "vis:oval")
    return shape(sl, MSO_SHAPE.OVAL, x, y, w, h, **kw)


def tri(sl, x, y, w, h, **kw):
    kw.setdefault("name", "vis:tri")
    return shape(sl, MSO_SHAPE.ISOSCELES_TRIANGLE, x, y, w, h, **kw)


def line(sl, x1, y1, x2, y2, color=None, lw=None, arrow=False,
         name="vis:line", dashed=False):
    cn = sl.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                 Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    cn.name = name
    _clean(cn)
    cn.line.color.rgb = color if color is not None else T.GREY_LT
    cn.line.width = lw or T.LINE_W
    ln = cn.line._get_or_add_ln()
    if arrow:
        te = ln.makeelement(qn("a:tailEnd"), {"type": "triangle",
                                              "w": "med", "len": "med"})
        ln.append(te)
    if dashed:
        d = ln.makeelement(qn("a:prstDash"), {"val": "dash"})
        ln.append(d)
    return cn


def freeform(sl, pts, fill=None, line_c=None, lw=None, name="vis:free",
             close=True):
    b = sl.shapes.build_freeform(Inches(pts[0][0]), Inches(pts[0][1]))
    b.add_line_segments([(Inches(px), Inches(py)) for px, py in pts[1:]],
                        close=close)
    sh = b.convert_to_shape()
    sh.name = name
    _clean(sh)
    if fill is None:
        _no_fill(sh)
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line_c is None:
        _no_line(sh)
    else:
        sh.line.color.rgb = line_c
        sh.line.width = lw or T.LINE_W
    return sh


# ==========================================================================
# Text
# ==========================================================================
def tbox(sl, x, y, w, h, paras, anchor="t", name="body:text", shrink=True,
         min_pt=T.MIN_PT):
    """
    Text box with zero internal margin and explicit line spacing, so the
    rendered height is exactly predictable (see textfit.py).
    """
    paras = [dict(p) for p in paras]
    if shrink:
        paras = autofit(paras, w, h, min_pt)

    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.name = name
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = _ANCHOR[anchor]

    for i, p in enumerate(paras):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = _ALIGN[p["align"]]
        para.line_spacing = Pt(p["pt"] * TF.LINE_FACTOR)
        para.space_after = Pt(p["space_after_pt"])
        para.space_before = Pt(0)
        if p["bullet_indent_in"]:
            para.paragraph_format.left_indent = Inches(p["bullet_indent_in"]) \
                if hasattr(para, "paragraph_format") else None
            pPr = para._p.get_or_add_pPr()
            pPr.set("marL", str(Emu(Inches(p["bullet_indent_in"]))))
            pPr.set("indent", "0")
        run = para.add_run()
        run.text = p["text"]
        f = run.font
        f.name = p["font"]
        f.size = Pt(p["pt"])
        f.bold = p["bold"]
        f.color.rgb = p["color"]
        if p["spc_pt"]:
            run._r.get_or_add_rPr().set("spc", str(int(p["spc_pt"] * 100)))
    return tb


def autofit(paras, w, h, min_pt=T.MIN_PT):
    """Scale the whole block down in 0.5pt steps until it fits. Never below min_pt."""
    scale = 1.0
    while scale > 0.5:
        trial = [dict(p, pt=max(min_pt, round(p["pt"] * scale * 2) / 2))
                 for p in paras]
        if TF.block_h_in(trial, w) <= h and TF.longest_line_in(trial, w) <= w + 0.002:
            return trial
        if all(p["pt"] * scale <= min_pt for p in paras):
            return [dict(p, pt=min_pt) for p in paras]
        scale -= 0.02
    return [dict(p, pt=min_pt) for p in paras]


def label_in(sh, text, font=None, pt=None, color=None, align="c", anchor="m",
             pad=0.10, bold=False, sa=0):
    """Put text inside an existing autoshape, autofit to the shape box."""
    w = sh.width / 914400.0 - 2 * pad
    h = sh.height / 914400.0 - 2 * pad
    paras = [P(text, font, pt, color, align=align, bold=bold, sa=sa)]
    paras = autofit(paras, w, h)
    tf = sh.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(pad)
    tf.margin_top = tf.margin_bottom = Inches(max(0.02, pad / 2))
    tf.vertical_anchor = _ANCHOR[anchor]
    para = tf.paragraphs[0]
    para.alignment = _ALIGN[align]
    para.line_spacing = Pt(paras[0]["pt"] * TF.LINE_FACTOR)
    run = para.add_run()
    run.text = text
    run.font.name = paras[0]["font"]
    run.font.size = Pt(paras[0]["pt"])
    run.font.bold = bold
    run.font.color.rgb = paras[0]["color"]
    return sh


# ==========================================================================
# Icon system — one line weight, drawn in code, never imported
# ==========================================================================
def icon(sl, kind, x, y, s=0.44, c=None):
    """Draw a line icon inside an s x s box at (x, y)."""
    c = c if c is not None else T.INK
    g = []
    u = s  # unit

    def L(x1, y1, x2, y2, arrow=False):
        g.append(line(sl, x + x1 * u, y + y1 * u, x + x2 * u, y + y2 * u,
                      c, T.LINE_W, arrow=arrow, name="vis:icon"))

    def O(cxp, cyp, r, fill=None):
        g.append(oval(sl, x + (cxp - r) * u, y + (cyp - r) * u, 2 * r * u, 2 * r * u,
                      fill=fill, line=None if fill else c, lw=T.LINE_W,
                      name="vis:icon"))

    def R(x1, y1, x2, y2, rounded=True, fill=None):
        f = rrect if rounded else rect
        g.append(f(sl, x + x1 * u, y + y1 * u, (x2 - x1) * u, (y2 - y1) * u,
                   fill=fill, line=None if fill else c, lw=T.LINE_W,
                   radius=0.05, name="vis:icon"))

    if kind == "model":                      # a small network = "the model"
        O(0.22, 0.28, 0.10); O(0.22, 0.74, 0.10); O(0.78, 0.51, 0.11)
        L(0.32, 0.31, 0.67, 0.48); L(0.32, 0.71, 0.67, 0.55)
    elif kind == "search":
        O(0.44, 0.42, 0.30); L(0.66, 0.64, 0.90, 0.88)
    elif kind == "chat":
        R(0.08, 0.14, 0.92, 0.70)
        L(0.28, 0.70, 0.24, 0.92); L(0.24, 0.92, 0.46, 0.70)
    elif kind == "doc":
        R(0.18, 0.06, 0.82, 0.94, rounded=True)
        L(0.32, 0.34, 0.68, 0.34); L(0.32, 0.50, 0.68, 0.50); L(0.32, 0.66, 0.56, 0.66)
    elif kind == "warn":
        g.append(freeform(sl, [(x + 0.5 * u, y + 0.06 * u),
                               (x + 0.96 * u, y + 0.90 * u),
                               (x + 0.04 * u, y + 0.90 * u)],
                          fill=None, line_c=c, lw=T.LINE_W, name="vis:icon"))
        L(0.5, 0.38, 0.5, 0.63); O(0.5, 0.77, 0.05, fill=c)
    elif kind == "lock":
        R(0.14, 0.44, 0.86, 0.94)
        g.append(shape(sl, MSO_SHAPE.ARC, x + 0.27 * u, y + 0.10 * u,
                       0.46 * u, 0.62 * u, fill=None, line=c, lw=T.LINE_W,
                       name="vis:icon"))
    elif kind == "shield":
        g.append(freeform(sl, [(x + 0.5 * u, y + 0.05 * u),
                               (x + 0.93 * u, y + 0.24 * u),
                               (x + 0.86 * u, y + 0.70 * u),
                               (x + 0.5 * u, y + 0.95 * u),
                               (x + 0.14 * u, y + 0.70 * u),
                               (x + 0.07 * u, y + 0.24 * u)],
                          fill=None, line_c=c, lw=T.LINE_W, name="vis:icon"))
    elif kind == "check":
        L(0.14, 0.52, 0.40, 0.78); L(0.40, 0.78, 0.88, 0.22)
    elif kind == "cross":
        L(0.18, 0.18, 0.82, 0.82); L(0.82, 0.18, 0.18, 0.82)
    elif kind == "person":
        O(0.5, 0.26, 0.19)
        g.append(freeform(sl, [(x + 0.10 * u, y + 0.95 * u),
                               (x + 0.16 * u, y + 0.60 * u),
                               (x + 0.84 * u, y + 0.60 * u),
                               (x + 0.90 * u, y + 0.95 * u)],
                          fill=None, line_c=c, lw=T.LINE_W, name="vis:icon",
                          close=False))
    elif kind == "clock":
        O(0.5, 0.5, 0.44); L(0.5, 0.5, 0.5, 0.24); L(0.5, 0.5, 0.72, 0.58)
    elif kind == "sheet":
        R(0.06, 0.12, 0.94, 0.88, rounded=True)
        L(0.06, 0.36, 0.94, 0.36); L(0.40, 0.12, 0.40, 0.88); L(0.06, 0.62, 0.94, 0.62)
    elif kind == "mail":
        R(0.06, 0.20, 0.94, 0.80, rounded=True)
        L(0.06, 0.24, 0.50, 0.56); L(0.50, 0.56, 0.94, 0.24)
    elif kind == "eye":
        g.append(freeform(sl, [(x + 0.04 * u, y + 0.50 * u),
                               (x + 0.50 * u, y + 0.14 * u),
                               (x + 0.96 * u, y + 0.50 * u),
                               (x + 0.50 * u, y + 0.86 * u)],
                          fill=None, line_c=c, lw=T.LINE_W, name="vis:icon"))
        O(0.5, 0.5, 0.14)
    elif kind == "key":
        O(0.26, 0.34, 0.22); L(0.40, 0.50, 0.92, 0.90)
        L(0.72, 0.70, 0.62, 0.82); L(0.84, 0.82, 0.74, 0.94)
    elif kind == "bulb":
        O(0.5, 0.36, 0.30); L(0.34, 0.74, 0.66, 0.74); L(0.38, 0.90, 0.62, 0.90)
    elif kind == "cycle":
        # 300-degree arc drawn as a polyline plus an arrow head — "do it again".
        # No masking rectangle, so it works on any background.
        arc = []
        for k in range(17):
            a = math.radians(-60 + k * (300.0 / 16))
            arc.append((x + (0.5 + 0.40 * math.cos(a)) * u,
                        y + (0.5 + 0.40 * math.sin(a)) * u))
        g.append(freeform(sl, arc, fill=None, line_c=c, lw=T.LINE_W,
                          name="vis:icon", close=False))
        g.append(freeform(sl, [(x + 0.439 * u, y + 0.074 * u),
                               (x + 0.345 * u, y + 0.232 * u),
                               (x + 0.255 * u, y + 0.076 * u)],
                          fill=c, name="vis:icon"))
    elif kind == "list":
        for i, yy in enumerate((0.20, 0.50, 0.80)):
            O(0.14, yy, 0.07, fill=c)
            L(0.32, yy, 0.94, yy)
    elif kind == "ban":
        O(0.5, 0.5, 0.44); L(0.19, 0.19, 0.81, 0.81)
    elif kind == "clip":
        R(0.24, 0.10, 0.76, 0.94)
        L(0.36, 0.34, 0.64, 0.34); L(0.36, 0.52, 0.64, 0.52); L(0.36, 0.70, 0.56, 0.70)
    else:
        O(0.5, 0.5, 0.44)
    return g


def icon_badge(sl, kind, x, y, size=0.78, bg=None, fg=None):
    """Icon inside a soft rounded tile — the standard 'section marker'."""
    bg = bg if bg is not None else T.SURFACE_ALT
    fg = fg if fg is not None else T.INK
    rrect(sl, x, y, size, size, fill=bg, radius=0.10, name="vis:badge")
    icon(sl, kind, x + size * 0.24, y + size * 0.24, size * 0.52, fg)


# ==========================================================================
# Diagram: numbered horizontal flow (process)
# ==========================================================================
def flow(sl, x, y, w, steps, ac, h=1.95, numbered=True):
    n = len(steps)
    gap = 0.42
    cw_ = (w - gap * (n - 1)) / n
    for i, st in enumerate(steps):
        sx = x + i * (cw_ + gap)
        card = rrect(sl, sx, y, cw_, h, fill=T.SURFACE_ALT, radius=0.10,
                     name="vis:flowcard")
        if numbered:
            oval(sl, sx + 0.20, y - 0.20, 0.42, 0.42, fill=ac, name="vis:flownum")
            tbox(sl, sx + 0.20, y - 0.115, 0.42, 0.26,
                 [P(str(i + 1), T.F_HEAD, 14, T.SURFACE, align="c")],
                 name="vis:flownumtext", shrink=False)
        tbox(sl, sx + 0.22, y + 0.40, cw_ - 0.44, 0.56,
             [P(st[0], T.F_HEAD, T.SZ_NODE + 2, T.INK)], name="vis:flowtitle")
        if len(st) > 1 and st[1]:
            tbox(sl, sx + 0.22, y + 1.00, cw_ - 0.44, h - 1.18,
                 [P(st[1], T.F_BODY, T.SZ_CAPTION, T.GREY)], name="vis:flowsub")
        if i < n - 1:
            line(sl, sx + cw_ + 0.09, y + h / 2, sx + cw_ + gap - 0.09, y + h / 2,
                 T.GREY_LT if False else ac, T.LINE_W, arrow=True, name="vis:flowarrow")
    return


# ==========================================================================
# Diagram: split-screen comparison
# ==========================================================================
def split(sl, x, y, w, h, left, right, ac, gap=0.34):
    """left/right: {'tag','title','items':[...],'tone': 'bad'|'good'|'accent'}"""
    pw = (w - gap) / 2
    for i, side in enumerate((left, right)):
        sx = x + i * (pw + gap)
        tone = side.get("tone", "accent")
        col = {"bad": T.ALERT, "good": T.SUCCESS, "accent": ac,
               "neutral": T.INK_SOFT}[tone]
        rrect(sl, sx, y, pw, h, fill=T.SURFACE_ALT, radius=0.10, name="vis:panel")
        rect(sl, sx, y, 0.075, h, fill=col, name="vis:panelspine")
        cy = y + 0.26
        chip(sl, sx + 0.30, cy, side["tag"], col)
        cy += 0.46
        tb = tbox(sl, sx + 0.30, cy, pw - 0.60, 0.62,
                  [P(side["title"], T.F_HEAD, T.SZ_NODE + 3, T.INK)],
                  name="vis:paneltitle")
        cy += 0.70
        items = side["items"]
        avail = (y + h - 0.24) - cy
        each = avail / max(1, len(items))
        for it in items:
            mark = side.get("mark") or ("check" if tone == "good"
                                        else ("cross" if tone == "bad" else "list"))
            icon(sl, mark, sx + 0.31, cy + 0.045, 0.20, col)
            tbox(sl, sx + 0.63, cy, pw - 0.93, each - 0.06,
                 [P(it, T.F_BODY, T.SZ_CAPTION + 1, T.INK_SOFT)], name="vis:panelitem")
            cy += each
    return


CHIP_PAD = 0.56          # total horizontal padding
CHIP_SLACK = 0.22        # extra room so PowerPoint never wraps a chip
CHIP_TXT_H = 0.24


def chip_w(text):
    return TF.text_w_in(text.upper(), T.F_HEAD, T.SZ_CHIP, 0.9) + CHIP_PAD


def chip(sl, x, y, text, col, h=0.30):
    w = chip_w(text)
    sh = rrect(sl, x, y, w, h, fill=col, radius=0.05, name="vis:chip")
    tbox(sl, x + (CHIP_PAD - CHIP_SLACK) / 2, y + (h - CHIP_TXT_H) / 2,
         w - CHIP_PAD + CHIP_SLACK, CHIP_TXT_H,
         [P(text.upper(), T.F_HEAD, T.SZ_CHIP, T.on(col), align="l", spc=0.9)],
         name="vis:chiptext", shrink=False)
    return sh, w


def chip_outline(sl, x, y, text, col, h=0.30):
    w = chip_w(text)
    rrect(sl, x, y, w, h, fill=None, line=col, lw=T.LINE_W_THIN, radius=0.05,
          name="vis:chip")
    tbox(sl, x + (CHIP_PAD - CHIP_SLACK) / 2, y + (h - CHIP_TXT_H) / 2,
         w - CHIP_PAD + CHIP_SLACK, CHIP_TXT_H,
         [P(text.upper(), T.F_HEAD, T.SZ_CHIP, col, align="l", spc=0.9)],
         name="vis:chiptext", shrink=False)
    return w


# ==========================================================================
# Diagram: nested block hierarchy
# ==========================================================================
NEST_HX = 0.36        # left/right inset per level
NEST_TOP = 0.94       # top inset — must clear the parent's label + sub line
NEST_BOT = 0.30


def nested(sl, x, y, w, h, layers, ac):
    """
    Concentric blocks, outermost first. Each child is inset far enough down
    that the parent's own label and caption stay fully visible.
    """
    n = len(layers)
    for i, lay in enumerate(layers):
        bx = x + i * NEST_HX
        by = y + i * NEST_TOP
        bw = w - 2 * i * NEST_HX
        bh = h - i * (NEST_TOP + NEST_BOT)
        last = (i == n - 1)
        shade = T.SURFACE if i % 2 else T.SURFACE_ALT
        rrect(sl, bx, by, bw, bh, fill=shade,
              line=ac if last else T.GREY_LT,
              lw=T.LINE_W if last else T.LINE_W_THIN, radius=0.10,
              name="vis:nest")
        tbox(sl, bx + 0.24, by + 0.16, bw - 0.48, 0.30,
             [P(lay["label"], T.F_HEAD, T.SZ_NODE + 1, ac if last else T.INK)],
             name="vis:nestlabel", shrink=False)
        if lay.get("sub"):
            tbox(sl, bx + 0.24, by + 0.50, bw - 0.48, 0.32,
                 [P(lay["sub"], T.F_BODY, T.SZ_CAPTION, T.GREY)],
                 name="vis:nestsub", shrink=False)
    return


def nested_min_h(n_layers, inner=1.10):
    return inner + (n_layers - 1) * (NEST_TOP + NEST_BOT)


# ==========================================================================
# Diagram: decision tree with explicit yes/no paths
# ==========================================================================
def tree(sl, x, y, w, question, yes, no, ac, h=3.30):
    """yes/no: {'label','detail','tone'}"""
    qw = w * 0.42
    qx = x + (w - qw) / 2
    tight = h < 3.70
    qh = 0.78 if tight else 0.92
    gap1 = 0.36 if tight else 0.46
    gap2 = 0.36 if tight else 0.46
    rrect(sl, qx, y, qw, qh, fill=T.INK, radius=0.10, name="vis:treeq")
    tbox(sl, qx + 0.22, y + 0.14, qw - 0.44, qh - 0.28,
         [P(question, T.F_HEAD, T.SZ_NODE + 1, T.SURFACE, align="c")],
         anchor="m", name="vis:treeqtext")

    midy = y + qh + gap1
    lx = x + w * 0.18
    rx = x + w * 0.82
    line(sl, x + w / 2, y + qh, x + w / 2, midy - 0.20, T.GREY_LT, T.LINE_W)
    line(sl, lx, midy - 0.20, rx, midy - 0.20, T.GREY_LT, T.LINE_W)

    bw = w * 0.40
    TONE_COL = {"good": T.SUCCESS, "bad": T.ALERT, "neutral": T.INK_SOFT}
    TONE_ICON = {"good": "check", "bad": "ban", "neutral": "search"}
    for side, cxp, tone_default in ((no, lx, "bad"), (yes, rx, "good")):
        tone = side.get("tone", tone_default)
        col = TONE_COL[tone]
        bx = cxp - bw / 2
        line(sl, cxp, midy - 0.20, cxp, midy + 0.10, T.GREY_LT, T.LINE_W, arrow=True)
        chip(sl, cxp - chip_w(side["path"]) / 2, midy - 0.05, side["path"], col)
        by = midy + gap2
        bh = max(1.55, y + h - by)
        rrect(sl, bx, by, bw, bh, fill=T.SURFACE_ALT, radius=0.10, name="vis:treebox")
        rect(sl, bx, by, bw, 0.075, fill=col, name="vis:treebar")
        icon(sl, TONE_ICON[tone], bx + 0.28, by + 0.34, 0.28, col)
        tbox(sl, bx + 0.68, by + 0.28, bw - 0.94, 0.40,
             [P(side["label"], T.F_HEAD, T.SZ_NODE + 3, T.INK)], name="vis:treelabel")
        tbox(sl, bx + 0.28, by + 0.80, bw - 0.56, bh - 1.02,
             [P(side["detail"], T.F_BODY, T.SZ_CAPTION + 3, T.INK_SOFT)],
             name="vis:treedetail")
    return


# ==========================================================================
# Diagram: annotated interface mock with red callout markers
# ==========================================================================
def mock_window(sl, x, y, w, h, title, rows, callouts=(), ac=None):
    """
    rows: list of (kind, text) where kind in {'field','text','muted','danger'}
    callouts: list of (row_index, note)
    """
    ac = ac if ac is not None else T.INK
    rrect(sl, x, y, w, h, fill=T.SURFACE, line=T.GREY_LT, lw=T.LINE_W_THIN,
          radius=0.08, name="vis:mock")
    rect(sl, x, y, w, 0.42, fill=T.SURFACE_ALT, name="vis:mockbar")
    for i in range(3):
        oval(sl, x + 0.16 + i * 0.20, y + 0.155, 0.11, 0.11, fill=T.GREY_LT,
             name="vis:mockdot")
    tbox(sl, x + 0.86, y + 0.11, w - 1.10, 0.24,
         [P(title, T.F_HEAD, T.SZ_CAPTION, T.GREY)], name="vis:mocktitle",
         shrink=False)

    ry = y + 0.62
    row_h = 0.46
    marks = {c[0]: c[1] for c in callouts}
    for i, (kind, text) in enumerate(rows):
        if kind == "field":
            rrect(sl, x + 0.26, ry, w - 0.52, row_h - 0.10, fill=T.SURFACE_ALT,
                  line=T.GREY_LT, lw=T.LINE_W_THIN, radius=0.05, name="vis:mockfield")
            tbox(sl, x + 0.44, ry + 0.055, w - 0.98, row_h - 0.22,
                 [P(text, T.F_MONO, T.SZ_CAPTION, T.INK_SOFT)], name="vis:mockrow")
        elif kind == "danger":
            rrect(sl, x + 0.26, ry, w - 0.52, row_h - 0.10, fill=T.SURFACE,
                  line=T.ALERT, lw=T.LINE_W, radius=0.05, name="vis:mockfield")
            tbox(sl, x + 0.44, ry + 0.055, w - 0.98, row_h - 0.22,
                 [P(text, T.F_MONO, T.SZ_CAPTION, T.ALERT)], name="vis:mockrow")
        else:
            col = T.GREY if kind == "muted" else T.INK_SOFT
            tbox(sl, x + 0.26, ry + 0.05, w - 0.52, row_h - 0.14,
                 [P(text, T.F_BODY, T.SZ_CAPTION, col)], name="vis:mockrow")
        if i in marks:
            mx = x + w - 0.02
            oval(sl, mx - 0.16, ry + 0.02, 0.30, 0.30, fill=T.ALERT, name="vis:callout")
            tbox(sl, mx - 0.16, ry + 0.075, 0.30, 0.20,
                 [P(str(list(marks).index(i) + 1), T.F_HEAD, T.SZ_CHIP, T.SURFACE,
                    align="c")], name="vis:calloutnum", shrink=False)
        ry += row_h
    return


def callout_notes(sl, x, y, w, notes, col=None):
    col = col if col is not None else T.ALERT
    cy = y
    for i, n in enumerate(notes):
        oval(sl, x, cy + 0.02, 0.28, 0.28, fill=col, name="vis:callout")
        tbox(sl, x, cy + 0.07, 0.28, 0.20,
             [P(str(i + 1), T.F_HEAD, T.SZ_CHIP, T.SURFACE, align="c")],
             name="vis:calloutnum", shrink=False)
        tbox(sl, x + 0.42, cy, w - 0.42, 0.62,
             [P(n, T.F_BODY, T.SZ_CAPTION + 1, T.INK_SOFT)], name="vis:calloutnote")
        cy += 0.72
    return


# ==========================================================================
# Prompt card — chat-bubble / terminal styling with a copy affordance
# ==========================================================================
PROMPT_PT = T.SZ_CAPTION + 2
_COPY_TXT = "Select + copy"


def prompt_card_h(text, w, caption=None):
    """Height this card needs. Lets callers size it to its content."""
    tw = w - 0.60
    th = TF.block_h_in([P(text, T.F_MONO, PROMPT_PT)], tw)
    ch = TF.block_h_in([P(caption, T.F_BODY, T.SZ_CAPTION)], tw) if caption else 0
    return 0.60 + th + (ch + 0.18 if caption else 0.0) + 0.34


def prompt_card(sl, x, y, w, h, text, ac, caption=None, header="Copy this prompt"):
    rrect(sl, x, y, w, h, fill=T.INK, radius=0.10, name="vis:promptcard")
    rect(sl, x, y, w, 0.44, fill=T.INK_SOFT, name="vis:promptbar")
    icon(sl, "chat", x + 0.24, y + 0.10, 0.24, T.SURFACE)

    # copy affordance — width measured, never guessed
    tw_ = TF.text_w_in(_COPY_TXT, T.F_HEAD, T.SZ_CHIP)
    bw = 0.30 + 0.17 + 0.12 + tw_ + 0.22
    bx = x + w - bw - 0.20
    rrect(sl, bx, y + 0.08, bw, 0.28, fill=ac, radius=0.05, name="vis:copybtn")
    icon(sl, "clip", bx + 0.16, y + 0.135, 0.17, T.SURFACE)
    tbox(sl, bx + 0.45, y + 0.145, tw_ + 0.14, 0.22,
         [P(_COPY_TXT, T.F_HEAD, T.SZ_CHIP, T.SURFACE)],
         name="vis:copytext", shrink=False)

    tbox(sl, x + 0.60, y + 0.12, bx - x - 0.72, 0.22,
         [P(header, T.F_HEAD, T.SZ_CAPTION, T.SURFACE, spc=0.6)],
         name="vis:prompthdr", shrink=False)

    tw = w - 0.60
    ch = TF.block_h_in([P(caption, T.F_BODY, T.SZ_CAPTION)], tw) if caption else 0
    body_h = h - 0.60 - 0.34 - (ch + 0.18 if caption else 0.0)
    tbox(sl, x + 0.30, y + 0.60, tw, max(0.3, body_h),
         [P(text, T.F_MONO, PROMPT_PT, T.SURFACE)], name="card:prompt")
    if caption:
        tbox(sl, x + 0.30, y + h - 0.30 - ch, tw, ch + 0.06,
             [P(caption, T.F_BODY, T.SZ_CAPTION, T.GREY_LT)], name="vis:promptcap",
             shrink=False)
    return


def output_card(sl, x, y, w, h, title, lines, ac, tone="good"):
    col = {"good": T.SUCCESS, "bad": T.ALERT, "accent": ac}[tone]
    rrect(sl, x, y, w, h, fill=T.SURFACE_ALT, radius=0.10, name="vis:outcard")
    rect(sl, x, y, 0.075, h, fill=col, name="vis:outspine")
    chip(sl, x + 0.28, y + 0.22, title, col)
    paras = []
    for i, ln in enumerate(lines):
        paras.append(P(ln, T.F_BODY, T.SZ_CAPTION + 1, T.INK_SOFT, sa=6))
    tbox(sl, x + 0.28, y + 0.70, w - 0.56, h - 0.94, paras, name="vis:outbody")
    return


# ==========================================================================
# Full-bleed alert band
# ==========================================================================
def alert_band(sl, y, h, headline, sub=None, col=None):
    col = col if col is not None else T.ALERT
    rect(sl, 0, y, T.SLIDE_W, h, fill=col, name="vis:alertband")
    icon(sl, "warn", T.MARGIN, y + (h - 0.46) / 2, 0.46, T.SURFACE)
    tx = T.MARGIN + 0.74
    if sub:
        tbox(sl, tx, y + 0.20, T.SLIDE_W - tx - T.MARGIN, h - 0.40,
             [P(headline, T.F_HEAD, 20, T.SURFACE, sa=4),
              P(sub, T.F_BODY, T.SZ_CAPTION + 2, T.SURFACE)],
             anchor="m", name="body:alert")
    else:
        tbox(sl, tx, y + 0.20, T.SLIDE_W - tx - T.MARGIN, h - 0.40,
             [P(headline, T.F_HEAD, 20, T.SURFACE)], anchor="m", name="body:alert")
    return


# ==========================================================================
# Checklist tiles
# ==========================================================================
MAX_TILE_H = 1.55


def checklist(sl, x, y, w, h, items, ac, cols=2, mark="check", dot=None):
    """
    Tiles never grow past MAX_TILE_H, so a short one-row list does not become
    a band of empty boxes. `mark` carries the meaning: a tick for things to do,
    a ban for things never to do.
    """
    dot = dot if dot is not None else ac
    rows = (len(items) + cols - 1) // cols
    gap = 0.24
    tw = (w - gap * (cols - 1)) / cols
    th = min(MAX_TILE_H, (h - gap * (rows - 1)) / rows)
    used = rows * th + gap * (rows - 1)
    y0 = y + max(0.0, (h - used) / 2)
    for i, it in enumerate(items):
        r, c = divmod(i, cols)
        bx = x + c * (tw + gap)
        by = y0 + r * (th + gap)
        rrect(sl, bx, by, tw, th, fill=T.SURFACE_ALT, radius=0.10,
              name="vis:checktile")
        oval(sl, bx + 0.24, by + (th - 0.34) / 2, 0.34, 0.34, fill=dot,
             name="vis:checkdot")
        icon(sl, mark, bx + 0.315, by + (th - 0.34) / 2 + 0.075, 0.19, T.SURFACE)
        tbox(sl, bx + 0.72, by + 0.16, tw - 0.96, th - 0.32,
             [P(it, T.F_BODY, T.SZ_CAPTION + 2, T.INK_SOFT)], anchor="m",
             name="cell:checkitem")
    return


# ==========================================================================
# Numbered mistake rows
# ==========================================================================
def mistake_rows(sl, x, y, w, h, items, ac):
    """items: list of (mistake, consequence)"""
    n = len(items)
    gap = 0.18
    rh = (h - gap * (n - 1)) / n
    for i, (m, cq) in enumerate(items):
        by = y + i * (rh + gap)
        rrect(sl, x, by, w, rh, fill=T.SURFACE_ALT, radius=0.08, name="vis:mistake")
        rect(sl, x, by, 0.075, rh, fill=T.ALERT, name="vis:mistakebar")
        oval(sl, x + 0.26, by + (rh - 0.34) / 2, 0.34, 0.34, fill=T.ALERT,
             name="vis:mistakenum")
        tbox(sl, x + 0.26, by + (rh - 0.34) / 2 + 0.075, 0.34, 0.22,
             [P(str(i + 1), T.F_HEAD, T.SZ_CHIP, T.SURFACE, align="c")],
             name="vis:mistakenumtext", shrink=False)
        half = (w - 1.10) * 0.46
        tbox(sl, x + 0.74, by + 0.10, half, rh - 0.20,
             [P(m, T.F_HEAD, T.SZ_CAPTION + 2, T.INK)], anchor="m", name="cell:mistake")
        line(sl, x + 0.74 + half + 0.16, by + 0.16, x + 0.74 + half + 0.16,
             by + rh - 0.16, T.GREY_LT, T.LINE_W_THIN)
        tbox(sl, x + 0.74 + half + 0.34, by + 0.10, w - (0.74 + half + 0.34) - 0.24,
             rh - 0.20, [P(cq, T.F_BODY, T.SZ_CAPTION + 1, T.GREY)], anchor="m",
             name="cell:mistakewhy")
    return


# ==========================================================================
# Step list for "Do this now"
# ==========================================================================
def picture(sl, path, x, y, w, h, name="vis:photo", frame=None):
    """
    Place an image cropped to fill the box exactly — never stretched, never
    letterboxed. Content files give a path; missing files fail the build
    loudly rather than leaving a placeholder box.
    """
    from PIL import Image
    if not os.path.exists(path):
        raise FileNotFoundError("image not found: %s" % path)
    iw, ih = Image.open(path).size
    box_ar = w / float(h)
    img_ar = iw / float(ih)
    pic = sl.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
    pic.name = name
    if img_ar > box_ar:                     # too wide — trim the sides
        c = (1 - box_ar / img_ar) / 2.0
        pic.crop_left = c
        pic.crop_right = c
    elif img_ar < box_ar:                   # too tall — trim top and bottom
        c = (1 - img_ar / box_ar) / 2.0
        pic.crop_top = c
        pic.crop_bottom = c
    if frame is not None:
        rrect(sl, x, y, w, h, fill=None, line=frame, lw=T.LINE_W_THIN,
              radius=0.10, name="vis:photoframe")
    return pic


def photo_card(sl, path, x, y, w, h, ac, caption=None, tag=None):
    """An image with the deck's card treatment: accent spine and a caption."""
    cap_h = 0.0
    if caption:
        cap_h = TF.block_h_in([P(caption, T.F_BODY, T.SZ_CAPTION)],
                              w - 0.10) + 0.22
    picture(sl, path, x, y, w, h - cap_h, frame=T.GREY_LT)
    if tag:
        chip(sl, x + 0.24, y + 0.24, tag, ac)
    if caption:
        tbox(sl, x + 0.04, y + h - cap_h + 0.14, w - 0.10, cap_h - 0.16,
             [P(caption, T.F_BODY, T.SZ_CAPTION, T.GREY)], name="cell:photocap")
    return


def why_list(sl, x, y, w, h, items, ac):
    """Small ticked rows — a structured reason list, not a bullet wall."""
    n = len(items)
    gap = 0.16
    rh = (h - gap * (n - 1)) / n
    for i, it in enumerate(items):
        by = y + i * (rh + gap)
        icon(sl, "check", x + 0.02, by + 0.06, 0.22, ac)
        tbox(sl, x + 0.42, by, w - 0.42, rh,
             [P(it, T.F_BODY, T.SZ_CAPTION + 2, T.INK_SOFT)], name="cell:why")
    return


def steps_list(sl, x, y, w, h, items, ac):
    n = len(items)
    gap = 0.16
    rh = (h - gap * (n - 1)) / n
    for i, it in enumerate(items):
        by = y + i * (rh + gap)
        oval(sl, x, by + (rh - 0.38) / 2, 0.38, 0.38, fill=ac, name="vis:stepnum")
        tbox(sl, x, by + (rh - 0.38) / 2 + 0.085, 0.38, 0.24,
             [P(str(i + 1), T.F_HEAD, T.SZ_NODE, T.SURFACE, align="c")],
             name="vis:stepnumtext", shrink=False)
        if i < n - 1:
            line(sl, x + 0.19, by + (rh + 0.38) / 2, x + 0.19,
                 by + rh + gap + (rh - 0.38) / 2, T.GREY_LT, T.LINE_W_THIN)
        tbox(sl, x + 0.60, by, w - 0.60, rh,
             [P(it, T.F_BODY, T.SZ_CAPTION + 3, T.INK_SOFT)], anchor="m",
             name="cell:step")
    return


# ==========================================================================
# Charts — the only PNG fallback. 2x resolution, transparent, deck palette.
# ==========================================================================
def _hex(c):
    return "#%02X%02X%02X" % (c[0], c[1], c[2])


def bar_chart_png(path, labels, values, ac, note_colour=None, ymax=None):
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt

    fig, axp = plt.subplots(figsize=(6.4, 3.0), dpi=200)
    fig.patch.set_alpha(0.0)
    axp.patch.set_alpha(0.0)
    cols = [_hex(ac)] * len(values)
    if note_colour is not None:
        cols[-1] = _hex(note_colour)
    bars = axp.bar(labels, values, color=cols, width=0.55)
    for b, v in zip(bars, values):
        axp.text(b.get_x() + b.get_width() / 2, v + (max(values) * 0.03),
                 f"{v}", ha="center", va="bottom",
                 fontsize=13, fontname="Segoe UI Semibold", color=_hex(T.INK))
    for s in ("top", "right", "left"):
        axp.spines[s].set_visible(False)
    axp.spines["bottom"].set_color(_hex(T.GREY_LT))
    axp.tick_params(axis="x", length=0, labelsize=12, colors=_hex(T.INK))
    axp.set_yticks([])
    if ymax:
        axp.set_ylim(0, ymax)
    for lbl in axp.get_xticklabels():
        lbl.set_fontname("Segoe UI")
    fig.tight_layout(pad=0.4)
    fig.savefig(path, transparent=True, dpi=200)
    plt.close(fig)
    return path
