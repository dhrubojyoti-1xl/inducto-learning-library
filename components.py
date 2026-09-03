"""
Reusable slide builders. Content files describe WHAT is on a slide;
this module decides HOW it is drawn. Coordinates come only from theme.py.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE

import os

import theme as T
import visuals as V
from visuals import P
import textfit as TF

ROOT = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.join(ROOT, "assets")


def _asset(path):
    """
    Resolve an image path from a content file. Content says
    "assets/whatever.png"; this finds it, or fails loudly. There is never a
    placeholder box.
    """
    if os.path.isabs(path):
        cand = path
    else:
        cand = os.path.join(ROOT, path)
        if not os.path.exists(cand):
            cand = os.path.join(ASSETS, os.path.basename(path))
    if not os.path.exists(cand):
        raise FileNotFoundError(
            "Image referenced by a content file is missing: %s -- "
            "drop it into %s and rebuild." % (path, ASSETS))
    return cand


# ==========================================================================
# Deck
# ==========================================================================
class Deck:
    def __init__(self, spec):
        self.spec = spec
        self.area = spec["area"]
        self.ac = T.accent(self.area)
        self.ac_lt = T.accent_light(self.area)     # accent that clears AA on INK
        self.prs = Presentation()
        self.prs.slide_width = Inches(T.SLIDE_W)
        self.prs.slide_height = Inches(T.SLIDE_H)
        self.blank = self.prs.slide_layouts[6]
        self.anchors = {}
        self.pending = []          # (shape, anchor)
        self.ext_links = []        # urls used
        self.slides = []
        self.points = 0            # content points

    # -- slide plumbing ---------------------------------------------------
    def new(self, anchor=None, hidden=False, bg=None):
        sl = self.prs.slides.add_slide(self.blank)
        if bg is not None:
            V.rect(sl, 0, 0, T.SLIDE_W, T.SLIDE_H, fill=bg, name="vis:bg")
        if hidden:
            sl._element.set("show", "0")
        if anchor:
            self.anchors[anchor] = sl
        self.slides.append(sl)
        return sl

    def link(self, shape, anchor):
        self.pending.append((shape, anchor))

    def ext(self, shape, url):
        shape.click_action.hyperlink.address = url
        if url not in self.ext_links:
            self.ext_links.append(url)

    def file_link(self, shape, relpath):
        """Relative hyperlink to another deck. Not an HTTP link."""
        shape.click_action.hyperlink.address = relpath
        self.file_links = getattr(self, "file_links", [])
        if relpath not in self.file_links:
            self.file_links.append(relpath)

    def resolve(self):
        missing = []
        for sh, anchor in self.pending:
            target = self.anchors.get(anchor)
            if target is None:
                missing.append(anchor)
                continue
            sh.click_action.target_slide = target
        return missing

    def save(self, path):
        missing = self.resolve()
        if missing:
            raise ValueError("unresolved anchors: %s" % sorted(set(missing)))
        self.prs.save(path)
        return path

    # -- chrome -----------------------------------------------------------
    def chrome(self, sl, label, title, nav=True, rule=True, dark=False):
        ink = T.SURFACE if dark else T.INK
        lab_c = T.SURFACE if dark else self.ac
        rule_c = T.INK_SOFT if dark else T.GREY_LT
        foot_c = T.GREY_LT if dark else T.GREY

        if label:
            V.tbox(sl, T.MARGIN, T.SEC_LABEL_Y, T.cw(8), T.SEC_LABEL_H,
                   [P(label.upper(), T.F_HEAD, T.SZ_LABEL, lab_c, spc=1.4)],
                   name="seclabel", shrink=False)
        if title:
            V.tbox(sl, T.MARGIN, T.TITLE_Y, T.TITLE_W, T.TITLE_H,
                   [P(title, T.F_HEAD, T.SZ_TITLE, ink)], anchor="b", name="title")
        if rule:
            V.line(sl, T.MARGIN, T.RULE_Y, T.SLIDE_W - T.MARGIN, T.RULE_Y,
                   rule_c, T.LINE_W_THIN, name="vis:rule")
        self.footer(sl, foot_c, dark)
        if nav:
            self.nav_button(sl, dark)

    def footer(self, sl, col=None, dark=False):
        col = col or T.GREY
        s = self.spec
        V.tbox(sl, T.MARGIN, T.FOOTER_Y, T.cw(7), T.FOOTER_H,
               [P("%s  ·  %s" % (s["module_code"],
                                 s.get("footer_title", s["title"])),
                  T.F_BODY, T.SZ_CAPTION, col)], name="footer", shrink=False)
        V.tbox(sl, T.SLIDE_W - T.MARGIN - T.cw(4), T.FOOTER_Y, T.cw(4), T.FOOTER_H,
               [P("Inducto Learning Library", T.F_BODY, T.SZ_CAPTION, col,
                  align="r")], name="footer", shrink=False)

    def nav_button(self, sl, dark=False):
        col = T.SURFACE if dark else self.ac
        sh = V.rrect(sl, T.NAV_X, T.NAV_Y, T.NAV_W, T.NAV_H,
                     fill=None, line=col, lw=T.LINE_W_THIN, radius=0.05,
                     name="nav:menu")
        V.icon(sl, "list", T.NAV_X + 0.16, T.NAV_Y + 0.105, 0.155, col)
        V.tbox(sl, T.NAV_X + 0.42, T.NAV_Y + 0.085, T.NAV_W - 0.56, 0.20,
               [P("Menu", T.F_HEAD, T.SZ_CHIP, col)], name="nav:menutext",
               shrink=False)
        self.link(sh, "menu")
        return sh

    def glossary_link(self, sl, x, y, term):
        w = V.chip_outline(sl, x, y, term + "  ?", T.GREY)
        hit = V.rrect(sl, x, y, w, 0.30, fill=None, name="nav:gloss")
        self.link(hit, "glossary")
        return w


# ==========================================================================
# Cover art — native shape compositions, one per motif
# ==========================================================================
def cover_art(sl, motif, ac, x, y, w, h):
    if motif == "network":
        nodes = [(0.16, 0.22), (0.10, 0.62), (0.42, 0.42), (0.44, 0.82),
                 (0.72, 0.20), (0.76, 0.62), (0.96, 0.40)]
        edges = [(0, 2), (1, 2), (2, 3), (2, 4), (2, 5), (4, 6), (5, 6), (3, 5)]
        pts = [(x + px * w, y + py * h) for px, py in nodes]
        for a, b in edges:
            V.line(sl, pts[a][0], pts[a][1], pts[b][0], pts[b][1],
                   T.INK_SOFT, T.LINE_W_THIN, name="vis:coverline")
        for i, (px, py) in enumerate(pts):
            r = 0.30 if i == 2 else (0.20 if i in (4, 5) else 0.14)
            V.oval(sl, px - r / 2, py - r / 2, r, r,
                   fill=ac if i in (2, 4, 5) else T.SURFACE, name="vis:covernode")
    elif motif == "prompt":
        V.rrect(sl, x, y + h * 0.06, w * 0.86, h * 0.40, fill=T.INK_SOFT,
                radius=0.10, name="vis:coverbubble")
        for i, ww in enumerate((0.60, 0.42, 0.52)):
            V.rect(sl, x + 0.30, y + h * 0.14 + i * 0.28, w * 0.86 * ww, 0.11,
                   fill=T.GREY, name="vis:coverline")
        V.rrect(sl, x + w * 0.14, y + h * 0.56, w * 0.86, h * 0.38, fill=ac,
                radius=0.10, name="vis:coverbubble")
        for i, ww in enumerate((0.66, 0.48)):
            V.rect(sl, x + w * 0.14 + 0.30, y + h * 0.64 + i * 0.28,
                   w * 0.86 * ww, 0.11, fill=T.SURFACE, name="vis:coverline")
    elif motif == "flow":
        n = 4
        bw, gap = w * 0.19, w * 0.08
        for i in range(n):
            bx = x + i * (bw + gap)
            by = y + h * 0.34 + (0.18 if i % 2 else -0.18)
            V.rrect(sl, bx, by, bw, bw * 0.86,
                    fill=ac if i == n - 1 else T.INK_SOFT, radius=0.10,
                    name="vis:coverblock")
            if i < n - 1:
                V.line(sl, bx + bw + 0.04, by + bw * 0.43,
                       bx + bw + gap - 0.04, by + bw * 0.43 + (0.36 if i % 2 else -0.36),
                       T.GREY, T.LINE_W_THIN, arrow=True, name="vis:coverline")
    elif motif == "layers":
        for i in range(4):
            V.rrect(sl, x + i * 0.34, y + h * 0.62 - i * 0.62, w * 0.72,
                    h * 0.24, fill=ac if i == 3 else T.INK_SOFT, radius=0.10,
                    name="vis:coverlayer")
    elif motif == "shield":
        cxp, cyp = x + w / 2, y + h / 2
        s = min(w, h) * 0.92
        V.freeform(sl, [(cxp, cyp - s * 0.50), (cxp + s * 0.42, cyp - s * 0.30),
                        (cxp + s * 0.36, cyp + s * 0.18), (cxp, cyp + s * 0.50),
                        (cxp - s * 0.36, cyp + s * 0.18),
                        (cxp - s * 0.42, cyp - s * 0.30)],
                   fill=ac, name="vis:covershield")
        V.freeform(sl, [(cxp, cyp - s * 0.34), (cxp + s * 0.28, cyp - s * 0.20),
                        (cxp + s * 0.24, cyp + s * 0.12), (cxp, cyp + s * 0.34),
                        (cxp - s * 0.24, cyp + s * 0.12),
                        (cxp - s * 0.28, cyp - s * 0.20)],
                   fill=None, line_c=T.SURFACE, lw=T.LINE_W, name="vis:covershield")
        V.icon(sl, "ban", cxp - s * 0.14, cyp - s * 0.14, s * 0.28, T.SURFACE)
    return


# ==========================================================================
# Structural slides
# ==========================================================================
PANEL_W = 6.05          # opaque text panel on an image cover


def cover(dk):
    s = dk.spec
    sl = dk.new("cover", bg=T.INK)
    img = s.get("cover_image")

    if img:
        # Hero image bleeds off the right edge; the text sits on an opaque
        # panel so contrast stays measurable rather than depending on the photo.
        V.picture(sl, _asset(img), PANEL_W, 0, T.SLIDE_W - PANEL_W, T.SLIDE_H,
                  name="vis:coverphoto")
        V.rect(sl, 0, 0, PANEL_W, T.SLIDE_H, fill=T.INK, name="vis:coverpanel")
        V.rect(sl, PANEL_W, 0, 0.07, T.SLIDE_H, fill=dk.ac, name="vis:coveredge")
        V.rect(sl, 0, 0, 0.22, T.SLIDE_H, fill=dk.ac, name="vis:coverspine")
        x0 = 0.92
        tw = PANEL_W - x0 - 0.55
    else:
        V.rect(sl, 0, 0, 0.22, T.SLIDE_H, fill=dk.ac, name="vis:coverspine")
        cover_art(sl, s.get("motif", "network"), dk.ac, 8.30, 1.70, 4.20, 3.60)
        x0 = 1.05
        tw = T.cw(6)
    # A long unbreakable word (e.g. "Problem-Solving") can exceed the panel
    # at 44pt, so fit the title first and measure the fitted version.
    title_p = V.autofit([P(s["title"], T.F_HEAD, T.SZ_COVER, T.SURFACE)],
                        tw, 2.60, min_pt=30)
    sub_p = V.autofit([P(s["subtitle"], T.F_BODY, T.SZ_SUB, T.GREY_LT)],
                      tw, 1.90, min_pt=16)
    th = TF.block_h_in(title_p, tw)
    sh_ = TF.block_h_in(sub_p, tw)

    # how many rows the chips will take, so the whole block centres correctly
    chip_rows, cx_ = 1, 0.0
    for text in (s["module_code"], "%d minutes" % s["duration_min"],
                 s["audience"]):
        need = V.chip_w(text)
        if cx_ > 0 and cx_ + need > tw:
            chip_rows += 1
            cx_ = 0.0
        cx_ += need + 0.16

    # Everything above the hint line. The hint itself sits at a fixed y so it
    # can never drift down into the footer.
    HINT_Y = T.FOOTER_Y - 0.62
    block = (0.30 + 0.20 + th + 0.36 + 0.30 + sh_ + 0.50
             + 0.34 + 0.46 * (chip_rows - 1))
    y = min(max(0.95, (T.SLIDE_H - block) / 2 - 0.34), HINT_Y - 0.40 - block)

    V.tbox(sl, x0, y, tw, 0.30,
           [P(T.AREAS[dk.area]["name"].upper(), T.F_HEAD, T.SZ_LABEL, dk.ac_lt,
              spc=1.4)], name="seclabel", shrink=False)
    y += 0.50
    V.tbox(sl, x0, y, tw, th + 0.06, title_p, name="title:cover",
           shrink=False)
    y += th + 0.36
    V.line(sl, x0, y, x0 + T.cw(4), y, T.INK_SOFT, T.LINE_W_THIN, name="vis:rule")
    y += 0.30
    V.tbox(sl, x0, y, tw, sh_ + 0.06, sub_p, name="body:coversub", shrink=False)
    y += sh_ + 0.50

    # Chips wrap onto a second row rather than running off the panel.
    x, rows = x0, 1
    for text, filled in ((s["module_code"], True),
                         ("%d minutes" % s["duration_min"], False),
                         (s["audience"], False)):
        need = V.chip_w(text)
        if x > x0 and x + need > x0 + tw:
            x = x0
            y += 0.46
            rows += 1
        if filled:
            _, w = V.chip(sl, x, y, text, dk.ac, h=0.34)
        else:
            w = V.chip_outline(sl, x, y, text, T.GREY_LT, h=0.34)
        x += w + 0.16

    y = HINT_Y
    V.tbox(sl, x0, y, tw, 0.48,
           [P("Start here  →  the next slide shows what this saves you.",
              T.F_BODY, T.SZ_CAPTION, T.GREY_DK)], name="cell:coverhint",
           shrink=False)
    dk.footer(sl, T.GREY_DK, dark=True)
    return sl


def why(dk):
    s = dk.spec["why"]
    sl = dk.new("why")
    dk.chrome(sl, "Why this matters to you", s["title"])

    lw = T.cw(7)
    top = T.BODY_TOP
    cost_h = 0.86
    story_h = T.BODY_H - cost_h - 0.26

    V.rrect(sl, T.MARGIN, top, lw, story_h, fill=T.SURFACE_ALT, radius=0.10,
            name="vis:story")
    V.rect(sl, T.MARGIN, top, 0.075, story_h, fill=dk.ac, name="vis:storyspine")
    V.icon_badge(sl, s.get("icon", "person"), T.MARGIN + 0.34, top + 0.34,
                 0.76, T.SURFACE, dk.ac)
    V.tbox(sl, T.MARGIN + 0.34, top + 1.34, lw - 0.68, story_h - 1.62,
           [P(s["scenario"], T.F_BODY, T.SZ_BODY, T.INK)], name="body:scenario")

    cy = top + story_h + 0.26
    V.rrect(sl, T.MARGIN, cy, lw, cost_h, fill=T.SURFACE, line=T.ALERT,
            lw=T.LINE_W_THIN, radius=0.08, name="vis:cost")
    V.rect(sl, T.MARGIN, cy, 0.075, cost_h, fill=T.ALERT, name="vis:costspine")
    V.icon(sl, "warn", T.MARGIN + 0.34, cy + (cost_h - 0.36) / 2, 0.36, T.ALERT)
    V.tbox(sl, T.MARGIN + 0.94, cy + 0.14, lw - 1.24, cost_h - 0.28,
           [P(s["cost"], T.F_HEAD, T.SZ_CAPTION + 4, T.ALERT)], anchor="m",
           name="body:cost")

    rx = T.cx(8)
    rw = T.cw(4)
    rh = T.BODY_H
    V.rrect(sl, rx, top, rw, rh, fill=T.INK, radius=0.10, name="vis:fix")
    V.tbox(sl, rx + 0.36, top + 0.38, rw - 0.72, 0.30,
           [P("AFTER THIS MODULE", T.F_HEAD, T.SZ_LABEL, dk.ac_lt, spc=1.4)],
           name="seclabel", shrink=False)
    V.line(sl, rx + 0.36, top + 0.86, rx + rw - 0.36, top + 0.86, T.INK_SOFT,
           T.LINE_W_THIN, name="vis:rule")
    V.tbox(sl, rx + 0.36, top + 1.10, rw - 0.72, rh - 1.50,
           [P(s["fix"], T.F_HEAD, 21, T.SURFACE)], name="body:fix")
    return sl


def outcomes(dk):
    s = dk.spec
    sl = dk.new("outcomes")
    dk.chrome(sl, "What you'll be able to do", "By the end, you can do these")
    items = s["outcomes"]
    n = len(items)
    gap = 0.20
    rh = (T.BODY_H - gap * (n - 1)) / n
    for i, (ico, text) in enumerate(items):
        y = T.BODY_TOP + i * (rh + gap)
        V.rrect(sl, T.MARGIN, y, T.CONTENT_W, rh, fill=T.SURFACE_ALT, radius=0.10,
                name="vis:outcome")
        V.icon_badge(sl, ico, T.MARGIN + 0.30, y + (rh - 0.62) / 2, 0.62,
                     T.SURFACE, dk.ac)
        V.tbox(sl, T.MARGIN + 1.14, y + 0.12, T.CONTENT_W - 1.50, rh - 0.24,
               [P(text, T.F_BODY, T.SZ_BODY, T.INK)], anchor="m",
               name="cell:outcome")
    return sl


def menu(dk, sections):
    sl = dk.new("menu")
    dk.chrome(sl, "Interactive menu", "Tap any line to jump straight there",
              nav=False)
    V.tbox(sl, T.MARGIN, T.RULE_Y + 0.06, T.CONTENT_W, 0.24,
           [P("Every slide has a Menu button in the top-right corner to bring you "
              "back here.", T.F_BODY, T.SZ_CAPTION, T.GREY)],
           name="cell:menuhint", shrink=False)

    rows = sections + [("Quiz", "5 questions, instant feedback", "quiz1"),
                       ("Recap card", "The whole topic on one screen", "recap"),
                       ("Your toolkit", "Templates, links, what to do next", "toolkit")]
    n = len(rows)
    gap = 0.10
    top = T.BODY_TOP + 0.12
    rh = (T.BODY_BOTTOM - top - gap * (n - 1)) / n
    for i, (label, sub, anchor) in enumerate(rows):
        y = top + i * (rh + gap)
        sh = V.rrect(sl, T.MARGIN, y, T.CONTENT_W, rh, fill=T.SURFACE_ALT,
                     radius=0.08, name="nav:menurow")
        V.rect(sl, T.MARGIN, y, 0.075, rh, fill=dk.ac, name="vis:menubar")
        V.tbox(sl, T.MARGIN + 0.34, y + 0.06, 0.44, rh - 0.12,
               [P("%02d" % (i + 1), T.F_HEAD, T.SZ_NODE, dk.ac)], anchor="m",
               name="vis:menunum", shrink=False)
        V.tbox(sl, T.MARGIN + 0.92, y + 0.06, T.cw(5), rh - 0.12,
               [P(label, T.F_HEAD, T.SZ_NODE + 3, T.INK)], anchor="m",
               name="cell:menulabel")
        V.tbox(sl, T.cx(7), y + 0.06, T.cw(4), rh - 0.12,
               [P(sub, T.F_BODY, T.SZ_CAPTION + 1, T.GREY)], anchor="m",
               name="cell:menusub")
        V.line(sl, T.SLIDE_W - T.MARGIN - 0.52, y + rh / 2,
               T.SLIDE_W - T.MARGIN - 0.28, y + rh / 2, dk.ac, T.LINE_W, arrow=True,
               name="vis:menuarrow")
        dk.link(sh, anchor)
    return sl


# ==========================================================================
# Content slide dispatcher
# ==========================================================================
def content_slide(dk, spec):
    sl = dk.new(spec.get("anchor"))
    dk.chrome(sl, spec.get("label", ""), spec["title"])

    y = T.BODY_TOP
    if spec.get("lead"):
        h = 0.30 + 0.30 * len(TF.wrap(spec["lead"], T.F_BODY, T.SZ_BODY, T.cw(9)))
        V.tbox(sl, T.MARGIN, y, T.cw(9), h,
               [P(spec["lead"], T.F_BODY, T.SZ_BODY, T.INK_SOFT)], name="body:lead")
        y += h + 0.18

    bottom = T.BODY_BOTTOM - (0.52 if spec.get("gloss") else 0.0)

    vis = spec.get("visual")
    if vis:
        draw_visual(dk, sl, vis, y, bottom - y)

    if spec.get("gloss"):
        gx = T.MARGIN
        V.tbox(sl, gx, bottom + 0.16, 1.76, 0.26,
               [P("What does this mean?", T.F_BODY, T.SZ_CAPTION, T.GREY)],
               name="cell:glosshint", shrink=False)
        gx += 1.88
        for term in spec["gloss"]:
            gx += dk.glossary_link(sl, gx, bottom + 0.13, term) + 0.14
    return sl


def draw_visual(dk, sl, vis, y, h):
    t = vis["type"]
    x, w = T.MARGIN, T.CONTENT_W
    ac = dk.ac

    if t == "flow":
        V.flow(sl, x, y + min(0.55, max(0.15, (h - 2.10) / 2)), w, vis["steps"], ac,
               h=min(2.10, h - 0.30))
    elif t == "split":
        V.split(sl, x, y, w, h, vis["left"], vis["right"], ac)
    elif t == "nested":
        nh = max(V.nested_min_h(len(vis["layers"])), min(h, 3.60))
        ny = y + max(0.0, (h - nh) / 2)
        V.nested(sl, x, ny, w * 0.56, nh, vis["layers"], ac)
        if vis.get("note"):
            V.tbox(sl, x + w * 0.62, ny + 0.16, w * 0.38, nh - 0.32,
                   [P(vis["note"], T.F_BODY, T.SZ_CAPTION + 3, T.GREY)],
                   name="body:nestnote")
    elif t == "tree":
        V.tree(sl, x, y + 0.10, w, vis["question"], vis["yes"], vis["no"], ac,
               h=h - 0.10)
    elif t == "mock":
        mw = w * 0.56
        V.mock_window(sl, x, y, mw, min(h, 3.90), vis["window"], vis["rows"],
                      vis.get("callouts", ()), ac)
        V.callout_notes(sl, x + mw + 0.40, y + 0.18, w - mw - 0.40, vis["notes"])
    elif t == "checklist":
        V.checklist(sl, x, y, w, h, vis["items"], ac, vis.get("cols", 2),
                    mark=vis.get("mark", "check"),
                    dot=T.ALERT if vis.get("mark") == "ban" else None)
    elif t == "mistakes":
        V.mistake_rows(sl, x, y, w, h, vis["items"], ac)
    elif t == "steps":
        sw = w * 0.50
        V.steps_list(sl, x, y, sw, h, vis["items"], ac)
        if vis.get("prompt"):
            pw = w - sw - 0.40
            ph = min(h, max(2.30, V.prompt_card_h(vis["prompt"], pw,
                                                  vis.get("caption"))))
            V.prompt_card(sl, x + sw + 0.40, y, pw, ph,
                          vis["prompt"], ac, vis.get("caption"),
                          vis.get("header", "Copy this prompt"))
    elif t == "prompt":
        pw = w * 0.62
        ph = min(h, max(2.30, V.prompt_card_h(vis["text"], pw, vis.get("caption"))))
        V.prompt_card(sl, x, y, pw, ph, vis["text"], ac,
                      vis.get("caption"), vis.get("header", "Copy this prompt"))
        if vis.get("why"):
            wx = x + pw + 0.40
            ww = w - pw - 0.40
            V.tbox(sl, wx, y + 0.06, ww, 0.28,
                   [P("WHY THIS WORKS", T.F_HEAD, T.SZ_LABEL, ac, spc=1.4)],
                   name="seclabel", shrink=False)
            V.why_list(sl, wx, y + 0.48, ww, min(h - 0.54, 2.90),
                       vis["why"], ac)
    elif t == "prompt_out":
        pw = w * 0.52
        ph = min(h, max(2.60, V.prompt_card_h(vis["text"], pw, vis.get("caption"))))
        V.prompt_card(sl, x, y, pw, ph, vis["text"], ac, vis.get("caption"))
        V.output_card(sl, x + pw + 0.40, y, w - pw - 0.40, ph,
                      vis.get("out_title", "What you get back"), vis["out"], ac,
                      "accent")
    elif t == "beforeafter":
        note_h = 0.62 if vis.get("note") else 0.0
        ch = min(h - note_h, 3.60)
        cw_ = (w - 0.40) / 2
        V.output_card(sl, x, y, cw_, ch, vis.get("bad_tag", "Before"),
                      vis["bad"], ac, "bad")
        V.output_card(sl, x + cw_ + 0.40, y, cw_, ch, vis.get("good_tag", "After"),
                      vis["good"], ac, "good")
        if vis.get("note"):
            V.tbox(sl, x, y + ch + 0.20, w, note_h - 0.24,
                   [P(vis["note"], T.F_HEAD, T.SZ_CAPTION + 3, T.INK)],
                   name="cell:banote")
    elif t == "iconrow":
        items = vis["items"]
        n = len(items)
        gap = 0.30
        tw = (w - gap * (n - 1)) / n
        th = min(h, 3.20)
        for i, it in enumerate(items):
            bx = x + i * (tw + gap)
            V.rrect(sl, bx, y, tw, th, fill=T.SURFACE_ALT, radius=0.10,
                    name="vis:tile")
            V.icon_badge(sl, it["icon"], bx + 0.30, y + 0.32, 0.72, T.SURFACE, ac)
            V.tbox(sl, bx + 0.30, y + 1.24, tw - 0.60, 0.62,
                   [P(it["label"], T.F_HEAD, T.SZ_NODE + 3, T.INK)],
                   name="cell:tilelabel")
            V.tbox(sl, bx + 0.30, y + 1.94, tw - 0.60, th - 2.20,
                   [P(it["sub"], T.F_BODY, T.SZ_CAPTION + 1, T.GREY)],
                   name="cell:tilesub")
    elif t == "bandlist":
        V.alert_band(sl, y, 1.05, vis["headline"], vis.get("sub"),
                     T.ALERT if vis.get("tone", "alert") == "alert" else ac)
        # tiles hug the band instead of floating in the middle of the slide
        cols = vis.get("cols", 2)
        rows = (len(vis["items"]) + cols - 1) // cols
        band_h = min(h - 1.35, rows * V.MAX_TILE_H + 0.24 * (rows - 1))
        V.checklist(sl, x, y + 1.35, w, band_h, vis["items"], ac, cols,
                    mark=vis.get("mark", "check"),
                    dot=T.ALERT if vis.get("mark") == "ban" else None)
    elif t == "chart":
        import os
        png = vis["png"]
        V.rrect(sl, x, y, w * 0.60, min(h, 3.30), fill=T.SURFACE_ALT, radius=0.10,
                name="vis:chartcard")
        sl.shapes.add_picture(png, Inches(x + 0.24), Inches(y + 0.24),
                              width=Inches(w * 0.60 - 0.48))
        V.tbox(sl, x + w * 0.60 + 0.40, y + 0.20, w * 0.40 - 0.40, h - 0.40,
               [P(vis["note"], T.F_BODY, T.SZ_CAPTION + 3, T.INK_SOFT)],
               name="body:chartnote")
    elif t == "image":
        # Image on one side, explanation on the other.
        iw = w * vis.get("split", 0.52)
        ih = min(h, vis.get("height", 3.60))
        side = vis.get("side", "right")
        ix = x + (w - iw) if side == "right" else x
        tx = x if side == "right" else x + iw + 0.40
        V.photo_card(sl, _asset(vis["path"]), ix, y, iw, ih, ac,
                     vis.get("caption"), vis.get("tag"))
        if vis.get("points"):
            V.tbox(sl, tx, y + 0.02, w - iw - 0.40, 0.28,
                   [P(vis.get("label", "WHAT TO LOOK AT"), T.F_HEAD,
                      T.SZ_LABEL, ac, spc=1.4)], name="seclabel", shrink=False)
            V.why_list(sl, tx, y + 0.46, w - iw - 0.40, min(ih - 0.52, 2.90),
                       vis["points"], ac)
    elif t == "image_band":
        ih = min(h, vis.get("height", 3.10))
        V.photo_card(sl, _asset(vis["path"]), x, y, w, ih, ac,
                     vis.get("caption"), vis.get("tag"))
    elif t == "image_compare":
        cwd = (w - 0.40) / 2
        ih = min(h, vis.get("height", 3.40))
        V.photo_card(sl, _asset(vis["bad_path"]), x, y, cwd, ih, ac,
                     vis.get("bad_caption"), vis.get("bad_tag", "Before"))
        V.photo_card(sl, _asset(vis["good_path"]), x + cwd + 0.40, y, cwd, ih,
                     ac, vis.get("good_caption"), vis.get("good_tag", "After"))
    else:
        raise ValueError("unknown visual type: %s" % t)


# ==========================================================================
# Quiz — branching, one feedback slide per wrong answer
# ==========================================================================
def quiz(dk, questions):
    """
    questions: list of dicts
      {'q':..., 'answers':[{'text':..,'ok':bool,'why':..}, x4], 'stem':...}
    """
    n = len(questions)
    for qi, q in enumerate(questions):
        anchor = "quiz%d" % (qi + 1)
        nxt = "quiz%d" % (qi + 2) if qi + 1 < n else "recap"
        sl = dk.new(anchor)
        dk.chrome(sl, "Quiz  ·  Question %d of %d" % (qi + 1, n), q["q"])

        y = T.BODY_TOP
        if q.get("stem"):
            V.rrect(sl, T.MARGIN, y, T.CONTENT_W, 0.78, fill=T.SURFACE_ALT,
                    radius=0.08, name="vis:stem")
            V.icon(sl, "chat", T.MARGIN + 0.28, y + 0.24, 0.30, dk.ac)
            V.tbox(sl, T.MARGIN + 0.76, y + 0.12, T.CONTENT_W - 1.04, 0.54,
                   [P(q["stem"], T.F_BODY, T.SZ_CAPTION + 3, T.INK_SOFT)],
                   anchor="m", name="cell:stem")
            y += 1.00

        gap = 0.20
        aw = (T.CONTENT_W - gap) / 2
        ah = (T.BODY_BOTTOM - y - gap) / 2
        for ai, a in enumerate(q["answers"]):
            r, c = divmod(ai, 2)
            bx = T.MARGIN + c * (aw + gap)
            by = y + r * (ah + gap)
            sh = V.rrect(sl, bx, by, aw, ah, fill=T.SURFACE_ALT,
                         line=T.GREY_LT, lw=T.LINE_W_THIN, radius=0.08,
                         name="quiz:answer")
            V.oval(sl, bx + 0.26, by + (ah - 0.40) / 2, 0.40, 0.40, fill=dk.ac,
                   name="vis:answerkey")
            V.tbox(sl, bx + 0.26, by + (ah - 0.40) / 2 + 0.085, 0.40, 0.25,
                   [P("ABCD"[ai], T.F_HEAD, T.SZ_NODE, T.SURFACE, align="c")],
                   name="vis:answerkeytext", shrink=False)
            V.tbox(sl, bx + 0.82, by + 0.14, aw - 1.10, ah - 0.28,
                   [P(a["text"], T.F_BODY, T.SZ_CAPTION + 3, T.INK)], anchor="m",
                   name="cell:answer")
            dk.link(sh, "%s_fb%d" % (anchor, ai))

        # feedback slides
        for ai, a in enumerate(q["answers"]):
            fb = dk.new("%s_fb%d" % (anchor, ai), hidden=True)
            ok = a["ok"]
            col = T.SUCCESS if ok else T.ALERT
            V.rect(fb, 0, 0, T.SLIDE_W, 1.30, fill=col, name="vis:fbband")
            V.icon(fb, "check" if ok else "cross", T.MARGIN, 0.42, 0.46, T.SURFACE)
            V.tbox(fb, T.MARGIN + 0.74, 0.40, T.cw(9), 0.56,
                   [P("Correct — here's why" if ok else "Not quite — here's what "
                      "actually happens", T.F_HEAD, 26, T.SURFACE)],
                   anchor="m", name="title:banner")
            V.tbox(fb, T.MARGIN, 1.62, T.cw(8), 0.30,
                   [P("YOU CHOSE  %s.  %s" % ("ABCD"[ai], a["text"][:70]),
                      T.F_HEAD, T.SZ_LABEL, T.GREY, spc=1.2)],
                   name="seclabel", shrink=False)
            V.tbox(fb, T.MARGIN, 2.10, T.cw(7), 2.60,
                   [P(a["why"], T.F_BODY, 20, T.INK)], name="body:fbwhy")

            V.rrect(fb, T.cx(8), 2.10, T.cw(4), 2.20, fill=T.SURFACE_ALT,
                    radius=0.10, name="vis:fbcard")
            V.tbox(fb, T.cx(8) + 0.30, 2.36, T.cw(4) - 0.60, 0.28,
                   [P("REMEMBER", T.F_HEAD, T.SZ_LABEL, dk.ac, spc=1.4)],
                   name="seclabel", shrink=False)
            V.tbox(fb, T.cx(8) + 0.30, 2.76, T.cw(4) - 0.60, 1.34,
                   [P(a.get("remember", q["remember"]), T.F_BODY,
                      T.SZ_CAPTION + 3, T.INK_SOFT)], name="cell:fbremember")

            btn = V.rrect(fb, T.MARGIN, 5.30, 3.10, 0.62, fill=dk.ac, radius=0.06,
                          name="nav:fbnext")
            V.tbox(fb, T.MARGIN + 0.30, 5.48, 2.50, 0.28,
                   [P("Next question  →" if ok and qi + 1 < n else
                      ("See the recap  →" if ok else "Try this question again  →"),
                      T.F_HEAD, T.SZ_NODE, T.SURFACE)], name="nav:fbnexttext",
                   shrink=False)
            dk.link(btn, nxt if ok else anchor)
            dk.footer(fb)
            dk.nav_button(fb, dark=True)
    return


# ==========================================================================
# Scenario branching
# ==========================================================================
def scenario(dk, sc):
    anchor = sc.get("anchor", "scenario")
    sl = dk.new(anchor)
    dk.chrome(sl, "Choose what you'd do", sc["title"])
    V.rrect(sl, T.MARGIN, T.BODY_TOP, T.CONTENT_W, 1.20, fill=T.INK, radius=0.10,
            name="vis:sccard")
    V.icon(sl, "person", T.MARGIN + 0.32, T.BODY_TOP + 0.34, 0.52, T.SURFACE)
    V.tbox(sl, T.MARGIN + 1.06, T.BODY_TOP + 0.20, T.CONTENT_W - 1.40, 0.80,
           [P(sc["situation"], T.F_BODY, T.SZ_BODY, T.SURFACE)], anchor="m",
           name="body:situation")

    y = T.BODY_TOP + 1.52
    n = len(sc["choices"])
    gap = 0.26
    cwd = (T.CONTENT_W - gap * (n - 1)) / n
    ch = T.BODY_BOTTOM - y
    for i, ch_ in enumerate(sc["choices"]):
        bx = T.MARGIN + i * (cwd + gap)
        sh = V.rrect(sl, bx, y, cwd, ch, fill=T.SURFACE_ALT, line=T.GREY_LT,
                     lw=T.LINE_W_THIN, radius=0.10, name="quiz:choice")
        V.chip(sl, bx + 0.28, y + 0.28, "Option %s" % "ABC"[i], dk.ac)
        V.tbox(sl, bx + 0.28, y + 0.88, cwd - 0.56, ch - 1.50,
               [P(ch_["text"], T.F_BODY, T.SZ_CAPTION + 3, T.INK)],
               name="cell:choice")
        V.line(sl, bx + 0.28, y + ch - 0.42, bx + 0.60, y + ch - 0.42, dk.ac,
               T.LINE_W, arrow=True, name="vis:choicearrow")
        V.tbox(sl, bx + 0.72, y + ch - 0.56, cwd - 1.00, 0.28,
               [P("See what happens", T.F_HEAD, T.SZ_CHIP, dk.ac)],
               name="cell:choicego", shrink=False)
        dk.link(sh, "%s_out%d" % (anchor, i))

    for i, ch_ in enumerate(sc["choices"]):
        fb = dk.new("%s_out%d" % (anchor, i), hidden=True)
        tone = ch_["tone"]
        col = {"good": T.SUCCESS, "bad": T.ALERT, "ok": dk.ac}[tone]
        V.rect(fb, 0, 0, T.SLIDE_W, 1.30, fill=col, name="vis:fbband")
        V.icon(fb, {"good": "check", "bad": "cross", "ok": "warn"}[tone],
               T.MARGIN, 0.42, 0.46, T.SURFACE)
        V.tbox(fb, T.MARGIN + 0.74, 0.40, T.cw(9), 0.56,
               [P(ch_["headline"], T.F_HEAD, 26, T.SURFACE)], anchor="m",
               name="title:banner")
        V.tbox(fb, T.MARGIN, 1.66, T.cw(7), 2.90,
               [P(ch_["consequence"], T.F_BODY, 20, T.INK)], name="body:consequence")
        V.rrect(fb, T.cx(8), 1.66, T.cw(4), 2.40, fill=T.SURFACE_ALT, radius=0.10,
                name="vis:fbcard")
        V.tbox(fb, T.cx(8) + 0.30, 1.92, T.cw(4) - 0.60, 0.28,
               [P("THE RULE", T.F_HEAD, T.SZ_LABEL, dk.ac, spc=1.4)],
               name="seclabel", shrink=False)
        V.tbox(fb, T.cx(8) + 0.30, 2.32, T.cw(4) - 0.60, 1.50,
               [P(ch_["rule"], T.F_BODY, T.SZ_CAPTION + 3, T.INK_SOFT)],
               name="cell:rule")
        btn = V.rrect(fb, T.MARGIN, 5.16, 3.30, 0.62, fill=dk.ac, radius=0.06,
                      name="nav:scback")
        V.tbox(fb, T.MARGIN + 0.30, 5.34, 2.70, 0.28,
               [P("Back to the situation  →", T.F_HEAD, T.SZ_NODE, T.SURFACE)],
               name="nav:scbacktext", shrink=False)
        dk.link(btn, anchor)
        dk.footer(fb)
        dk.nav_button(fb, dark=True)
    return


# ==========================================================================
# Recap + toolkit + glossary
# ==========================================================================
def recap(dk, rc):
    sl = dk.new("recap")
    dk.chrome(sl, "Recap card", rc["title"])
    V.tbox(sl, T.MARGIN, T.RULE_Y + 0.10, T.CONTENT_W, 0.28,
           [P("Screenshot this slide. It is the whole module in one screen.",
              T.F_BODY, T.SZ_CAPTION, T.GREY)], name="cell:recaphint", shrink=False)

    y = T.BODY_TOP + 0.20
    items = rc["points"]
    n = len(items)
    gap = 0.18
    cols = 2
    rows = (n + cols - 1) // cols
    band_h = 0.82
    tiles_h = T.BODY_BOTTOM - y - band_h - 0.16
    tw = (T.CONTENT_W - gap) / cols
    th = (tiles_h - gap * (rows - 1)) / rows
    for i, (head, sub) in enumerate(items):
        r, c = divmod(i, cols)
        bx = T.MARGIN + c * (tw + gap)
        by = y + r * (th + gap)
        V.rrect(sl, bx, by, tw, th, fill=T.SURFACE_ALT, radius=0.08,
                name="vis:recaptile")
        V.rect(sl, bx, by, 0.06, th, fill=dk.ac, name="vis:recapbar")
        V.tbox(sl, bx + 0.26, by + 0.12, tw - 0.52, 0.30,
               [P(head, T.F_HEAD, T.SZ_CAPTION + 3, T.INK)], name="cell:recaphead")
        V.tbox(sl, bx + 0.26, by + 0.46, tw - 0.52, th - 0.58,
               [P(sub, T.F_BODY, T.SZ_CAPTION + 1, T.GREY)], name="cell:recapsub")

    by = T.BODY_BOTTOM - band_h
    V.rrect(sl, T.MARGIN, by, T.CONTENT_W, T.BODY_BOTTOM - by, fill=T.INK,
            radius=0.10, name="vis:recapline")
    V.icon(sl, "bulb", T.MARGIN + 0.32, by + (T.BODY_BOTTOM - by - 0.44) / 2, 0.44,
           dk.ac_lt)
    V.tbox(sl, T.MARGIN + 1.00, by + 0.14, T.CONTENT_W - 1.34,
           T.BODY_BOTTOM - by - 0.28,
           [P(rc["oneliner"], T.F_HEAD, 20, T.SURFACE)], anchor="m",
           name="cell:recapone")
    return sl


def toolkit(dk, tk):
    sl = dk.new("toolkit")
    dk.chrome(sl, "Your toolkit", tk.get("title", "Take these with you"))

    y = T.BODY_TOP
    lw = T.cw(7)
    V.tbox(sl, T.MARGIN, y, lw, 0.30,
           [P("TEMPLATES YOU CAN REUSE TODAY", T.F_HEAD, T.SZ_LABEL, dk.ac, spc=1.4)],
           name="seclabel", shrink=False)
    ty = y + 0.44
    n = len(tk["templates"])
    gap = 0.16
    rh = (2.80 - gap * (n - 1)) / n
    for i, (ico, name_, sub) in enumerate(tk["templates"]):
        by = ty + i * (rh + gap)
        V.rrect(sl, T.MARGIN, by, lw, rh, fill=T.SURFACE_ALT, radius=0.08,
                name="vis:tk")
        V.icon_badge(sl, ico, T.MARGIN + 0.24, by + (rh - 0.56) / 2, 0.56,
                     T.SURFACE, dk.ac)
        V.tbox(sl, T.MARGIN + 1.00, by + 0.10, lw - 1.28, rh * 0.50,
               [P(name_, T.F_HEAD, T.SZ_CAPTION + 3, T.INK)], name="cell:tkname")
        V.tbox(sl, T.MARGIN + 1.00, by + rh * 0.50 + 0.02, lw - 1.28, rh * 0.44,
               [P(sub, T.F_BODY, T.SZ_CAPTION, T.GREY)], name="cell:tksub")

    rx, rw = T.cx(8), T.cw(4)
    V.tbox(sl, rx, y, rw, 0.30,
           [P("LINKS THAT WORK", T.F_HEAD, T.SZ_LABEL, dk.ac, spc=1.4)],
           name="seclabel", shrink=False)
    ly = y + 0.44
    for label, url in tk["links"]:
        sh = V.rrect(sl, rx, ly, rw, 0.62, fill=T.SURFACE_ALT, radius=0.06,
                     name="nav:link")
        V.icon(sl, "search", rx + 0.24, ly + 0.16, 0.30, dk.ac)
        V.tbox(sl, rx + 0.68, ly + 0.09, rw - 0.94, 0.24,
               [P(label, T.F_HEAD, T.SZ_CAPTION + 1, T.INK)], name="cell:linklabel",
               shrink=False)
        shown = TF.truncate(url.replace("https://", "").rstrip("/"),
                            T.F_BODY, T.SZ_CAPTION, rw - 0.96)
        V.tbox(sl, rx + 0.68, ly + 0.33, rw - 0.94, 0.22,
               [P(shown, T.F_BODY, T.SZ_CAPTION, dk.ac)],
               name="cell:linkurl", shrink=False)
        dk.ext(sh, url)
        ly += 0.74

    ny = ly + 0.16
    V.rrect(sl, rx, ny, rw, T.BODY_BOTTOM - ny, fill=T.INK, radius=0.08,
            name="vis:next")
    V.tbox(sl, rx + 0.28, ny + 0.20, rw - 0.56, 0.26,
           [P("DO THIS NEXT", T.F_HEAD, T.SZ_LABEL, dk.ac_lt, spc=1.4)],
           name="seclabel", shrink=False)
    V.tbox(sl, rx + 0.28, ny + 0.56, rw - 0.56, T.BODY_BOTTOM - ny - 0.80,
           [P(tk["next"], T.F_BODY, T.SZ_CAPTION + 3, T.SURFACE)],
           name="cell:next")
    return sl


def video_slide(dk, v):
    """
    One verified third-party video. Title, channel and runtime are the values
    YouTube itself returned — never the ones a source claimed.
    """
    sl = dk.new("video")
    dk.chrome(sl, "Watch this", v.get("heading", "The same ideas, explained out loud"))

    lw = T.cw(7)
    h = T.BODY_H
    card = V.rrect(sl, T.MARGIN, T.BODY_TOP, lw, h, fill=T.INK, radius=0.10,
                   name="nav:video")
    # play affordance
    pcx, pcy = T.MARGIN + 1.24, T.BODY_TOP + 1.30
    V.oval(sl, pcx - 0.52, pcy - 0.52, 1.04, 1.04, fill=dk.ac, name="vis:play")
    V.freeform(sl, [(pcx - 0.15, pcy - 0.25), (pcx + 0.27, pcy),
                    (pcx - 0.15, pcy + 0.25)], fill=T.SURFACE, name="vis:playtri")

    V.tbox(sl, T.MARGIN + 2.20, T.BODY_TOP + 0.56, lw - 2.54, 1.50,
           [P(v["title"], T.F_HEAD, 21, T.SURFACE)], name="body:videotitle")

    # Two chips on one row, always inside the card. A long channel name is
    # elided rather than allowed to overflow the panel.
    cy = T.BODY_TOP + 2.34
    x = T.MARGIN + 0.40
    right = T.MARGIN + lw - 0.40
    dur = v["duration"] + " min"
    max_chan = right - x - V.chip_w(dur) - 0.16
    chan = v["channel"]
    while len(chan) > 4 and V.chip_w(chan) > max_chan:
        chan = chan[:-2]
    if chan != v["channel"]:
        chan = chan.rstrip(" -|,") + "…"
    x += V.chip_outline(sl, x, cy, chan, T.GREY_DK, h=0.32) + 0.16
    V.chip_outline(sl, x, cy, dur, T.GREY_DK, h=0.32)

    V.line(sl, T.MARGIN + 0.40, cy + 0.62, T.MARGIN + lw - 0.40, cy + 0.62,
           T.INK_SOFT, T.LINE_W_THIN, name="vis:rule")
    V.tbox(sl, T.MARGIN + 0.40, cy + 0.86, lw - 0.80, 0.60,
           [P(v["note"], T.F_BODY, T.SZ_CAPTION + 2, T.GREY_DK)],
           name="body:videonote")

    btn = V.rrect(sl, T.MARGIN + 0.40, T.BODY_TOP + h - 0.86, 3.30, 0.54,
                  fill=dk.ac, radius=0.06, name="nav:videobtn")
    V.icon(sl, "search", T.MARGIN + 0.58, T.BODY_TOP + h - 0.74, 0.28, T.SURFACE)
    V.tbox(sl, T.MARGIN + 1.00, T.BODY_TOP + h - 0.72, 2.50, 0.26,
           [P("Click to open on YouTube", T.F_HEAD, T.SZ_NODE, T.SURFACE)],
           name="nav:videobtntext", shrink=False)
    dk.ext(card, v["url"])
    dk.ext(btn, v["url"])

    rx, rw = T.cx(8), T.cw(4)
    V.rrect(sl, rx, T.BODY_TOP, rw, h, fill=T.SURFACE_ALT, radius=0.10,
            name="vis:videohow")
    V.rect(sl, rx, T.BODY_TOP, 0.075, h, fill=dk.ac, name="vis:videospine")
    V.tbox(sl, rx + 0.34, T.BODY_TOP + 0.34, rw - 0.68, 0.28,
           [P("HOW TO USE IT", T.F_HEAD, T.SZ_LABEL, dk.ac, spc=1.4)],
           name="seclabel", shrink=False)
    V.why_list(sl, rx + 0.34, T.BODY_TOP + 0.82, rw - 0.68, h - 1.20,
               v["how"], dk.ac)
    return sl


def glossary(dk, terms):
    sl = dk.new("glossary", hidden=True)
    dk.chrome(sl, "Glossary", "Words people use around AI")
    n = len(terms)
    cols = 2
    rows = (n + cols - 1) // cols
    gap = 0.18
    tw = (T.CONTENT_W - gap) / cols
    th = (T.BODY_H - 0.80 - gap * (rows - 1)) / rows
    for i, (term, d) in enumerate(terms):
        r, c = divmod(i, cols)
        bx = T.MARGIN + c * (tw + gap)
        by = T.BODY_TOP + r * (th + gap)
        V.rrect(sl, bx, by, tw, th, fill=T.SURFACE_ALT, radius=0.08,
                name="vis:gloss")
        V.rect(sl, bx, by, 0.06, th, fill=dk.ac, name="vis:glossbar")
        V.tbox(sl, bx + 0.26, by + 0.12, tw - 0.52, 0.30,
               [P(term, T.F_HEAD, T.SZ_CAPTION + 3, T.INK)], name="cell:glossterm")
        V.tbox(sl, bx + 0.26, by + 0.46, tw - 0.52, th - 0.58,
               [P(d, T.F_BODY, T.SZ_CAPTION + 1, T.GREY)], name="cell:glossdef")
    btn = V.rrect(sl, T.MARGIN, T.BODY_BOTTOM - 0.56, 2.90, 0.56, fill=dk.ac,
                  radius=0.06, name="nav:glossback")
    V.tbox(sl, T.MARGIN + 0.28, T.BODY_BOTTOM - 0.40, 2.34, 0.26,
           [P("Back to the menu  →", T.F_HEAD, T.SZ_NODE, T.SURFACE)],
           name="nav:glossbacktext", shrink=False)
    dk.link(btn, "menu")
    return sl
